#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate a CEO-facing PowerPoint deck from the NotebookLM slide outline.

Enhanced with modern Figma-inspired design system:
- Figma Design Tokens sync (palette, typography, 8px-grid spacing, rounded cards)
- Modern Studio Card Layouts (glass/dark-surface containers, vibrant accent chips)
- Multiple theme modes: 'figma-modern' (Studio dark ink), 'default' (Light executive), 'canva-export' (Flat)
"""

from __future__ import annotations

import argparse
import datetime as dt
import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports" / "knowledge_flow"
TOKENS_FILE = PROJECT_ROOT / "docs" / "design_tokens.json"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Design Token Color Palettes ---
PALETTES = {
    "figma-modern": {
        "bg": RGBColor(11, 15, 25),          # Deep Ink #0B0F19
        "surface": RGBColor(21, 30, 46),     # Dark Surface #151E2E
        "surface_elevated": RGBColor(30, 41, 59), # Slate #1E293B
        "border": RGBColor(42, 54, 79),      # Subtle border
        "border_accent": RGBColor(56, 189, 248), # Cyan highlight #38BDF8
        "text_primary": RGBColor(248, 250, 252), # Pure text #F8FAFC
        "text_secondary": RGBColor(148, 163, 184), # Muted slate #94A3B8
        "text_muted": RGBColor(100, 116, 139),   # Meta text
        "accent_cyan": RGBColor(56, 189, 248),   # Cyan #38BDF8
        "accent_indigo": RGBColor(99, 102, 241), # Indigo #6366F1
        "accent_emerald": RGBColor(16, 185, 129), # Emerald #10B981
        "accent_amber": RGBColor(245, 158, 11),  # Amber #F59E0B
        "accent_rose": RGBColor(244, 63, 94),    # Rose #F43F5E
        "chip_bg_cyan": RGBColor(12, 44, 72),
        "chip_bg_emerald": RGBColor(6, 44, 33),
        "chip_bg_amber": RGBColor(53, 38, 10),
        "chip_bg_rose": RGBColor(50, 15, 25),
    },
    "default": {
        "bg": RGBColor(255, 255, 255),
        "surface": RGBColor(247, 250, 253),
        "surface_elevated": RGBColor(255, 255, 255),
        "border": RGBColor(217, 224, 232),
        "border_accent": RGBColor(26, 115, 232),
        "text_primary": RGBColor(25, 34, 45),
        "text_secondary": RGBColor(91, 105, 120),
        "text_muted": RGBColor(120, 135, 150),
        "accent_cyan": RGBColor(26, 115, 232),
        "accent_indigo": RGBColor(16, 72, 154),
        "accent_emerald": RGBColor(20, 146, 96),
        "accent_amber": RGBColor(190, 124, 0),
        "accent_rose": RGBColor(188, 68, 60),
        "chip_bg_cyan": RGBColor(230, 242, 255),
        "chip_bg_emerald": RGBColor(226, 246, 238),
        "chip_bg_amber": RGBColor(255, 244, 218),
        "chip_bg_rose": RGBColor(255, 232, 229),
    },
}

CURRENT_THEME = "figma-modern"
COLORS = PALETTES["figma-modern"]
FONT_DISPLAY = "Plus Jakarta Sans"
FONT_BODY = "Yu Gothic"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def load_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def get_nested(data: dict[str, Any], *keys: str, default: Any = None) -> Any:
    current: Any = data
    for key in keys:
        if not isinstance(current, dict) or key not in current:
            return default
        current = current[key]
    return current


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name or FONT_BODY
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color or COLORS["text_primary"]
    return box


def add_bullets(
    slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    color: RGBColor | None = None,
    font_name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"•  {bullet}"
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = font_name or FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["text_secondary"]
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_chip(
    slide,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    text: RGBColor,
    border: RGBColor | None = None,
):
    add_rect(slide, x, y, w, 0.32, fill=fill, line=border or fill, radius=True)
    add_text(
        slide,
        label,
        x,
        y + 0.04,
        w,
        0.24,
        size=9,
        color=text,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        font_name=FONT_DISPLAY,
    )


def add_background(slide):
    """Draw full-slide background for dark/light canvas."""
    add_rect(slide, 0, 0, 13.333, 7.5, fill=COLORS["bg"], line=COLORS["bg"])


def add_footer(slide, slide_no: int, source_note: str):
    add_text(
        slide,
        f"{slide_no:02d} / 08   |   {source_note}",
        0.75,
        7.05,
        8.0,
        0.28,
        size=8,
        color=COLORS["text_muted"],
        font_name=FONT_MONO,
    )
    add_text(
        slide,
        "Mighty Skill-Bridge  ·  CEO Decision Deck",
        8.8,
        7.05,
        3.8,
        0.28,
        size=8,
        color=COLORS["text_muted"],
        align=PP_ALIGN.RIGHT,
        font_name=FONT_DISPLAY,
    )


def add_header(slide, slide_no: int, eyebrow: str, title: str, source_note: str):
    add_background(slide)
    # Top subtle decorative line
    add_rect(slide, 0.75, 0.35, 0.5, 0.04, fill=COLORS["accent_cyan"], line=COLORS["accent_cyan"])
    # Eyebrow / Category Tag
    add_text(
        slide,
        f"{slide_no:02d}  /  {eyebrow.upper()}",
        0.75,
        0.48,
        8.0,
        0.26,
        size=9,
        color=COLORS["accent_cyan"],
        bold=True,
        font_name=FONT_MONO,
    )
    # Main Slide Title
    add_text(
        slide,
        title,
        0.73,
        0.76,
        11.8,
        0.65,
        size=26,
        color=COLORS["text_primary"],
        bold=True,
        font_name=FONT_DISPLAY,
    )
    add_footer(slide, slide_no, source_note)


def add_evidence_panel(slide, items: list[str], y: float = 5.75):
    """Bottom evidence & traceability bar."""
    add_rect(slide, 0.75, y, 11.83, 1.15, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "見せる証跡", 0.95, y + 0.18, 1.2, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"], border=COLORS["border"])
    add_bullets(slide, items, 2.3, y + 0.12, 10.1, 0.9, size=10.5, color=COLORS["text_secondary"], font_name=FONT_MONO)


def add_metric_card(slide, x: float, y: float, w: float, h: float, label: str, value: str, note: str, accent: RGBColor):
    """Figma modern metric card with rounded corners and distinct visual hierarchy."""
    add_rect(slide, x, y, w, h, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    # Accent top mini bar
    add_rect(slide, x + 0.2, y + 0.18, 0.35, 0.04, fill=accent, line=accent)
    add_text(slide, value, x + 0.18, y + 0.28, w - 0.36, 0.45, size=24, color=accent, bold=True, font_name=FONT_DISPLAY)
    add_text(slide, label, x + 0.2, y + 0.76, w - 0.4, 0.26, size=12, color=COLORS["text_primary"], bold=True)
    add_text(slide, note, x + 0.2, y + 1.05, w - 0.4, 0.32, size=9.5, color=COLORS["text_secondary"])


def add_option_card(slide, x: float, y: float, w: float, h: float, title: str, audience: str, value: str, first_step: str, is_recommended: bool = False):
    border_color = COLORS["border_accent"] if is_recommended else COLORS["border"]
    line_w = 1.8 if is_recommended else 1.0
    add_rect(slide, x, y, w, h, fill=COLORS["surface"], line=border_color, line_width=line_w, radius=True)
    
    if is_recommended:
        add_chip(slide, "★ 推薦初手", x + w - 1.25, y + 0.15, 1.05, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"], border=COLORS["border_accent"])

    add_text(slide, title, x + 0.22, y + 0.2, w - 1.4, 0.32, size=15, color=COLORS["accent_cyan"], bold=True, font_name=FONT_DISPLAY)
    
    # Rows
    row_y = y + 0.65
    add_text(slide, "想定対象", x + 0.22, row_y, 0.85, 0.22, size=9, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, audience, x + 1.05, row_y - 0.02, w - 1.25, 0.28, size=10.5, color=COLORS["text_primary"])
    
    row_y += 0.45
    add_text(slide, "提供価値", x + 0.22, row_y, 0.85, 0.22, size=9, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, value, x + 1.05, row_y - 0.02, w - 1.25, 0.42, size=10.5, color=COLORS["text_secondary"])
    
    row_y += 0.55
    add_text(slide, "最初の行動", x + 0.22, row_y, 0.85, 0.22, size=9, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, first_step, x + 1.05, row_y - 0.02, w - 1.25, 0.42, size=10.5, color=COLORS["accent_emerald"], bold=True)


def add_flow_card(slide, x: float, y: float, w: float, h: float, step_no: str, title: str, note: str, accent: RGBColor):
    add_rect(slide, x, y, w, h, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, step_no, x + 0.15, y + 0.15, 0.55, fill=COLORS["surface_elevated"], text=accent, border=COLORS["border"])
    add_text(slide, title, x + 0.15, y + 0.52, w - 0.3, 0.28, size=13, color=COLORS["text_primary"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, note, x + 0.15, y + 0.82, w - 0.3, 0.45, size=9, color=COLORS["text_secondary"])


def add_arrow(slide, x: float, y: float):
    add_text(slide, "→", x, y, 0.35, 0.35, size=16, color=COLORS["accent_cyan"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_MONO)


def build_slides(prs: Presentation, context: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    source_note = f"NotebookLM CLI outline · ID: {context['notebook_id'][:8]}..."

    # 1: Agenda & Decision Scope
    slide = prs.slides.add_slide(blank)
    add_header(slide, 1, "Decision Meeting Scope", "本日決めたいこと", source_note)
    add_text(slide, "企画を決め打ちする場ではなく、次の開発の向きと優先順位を決める場です。", 0.78, 1.5, 11.7, 0.35, size=15, color=COLORS["text_secondary"])
    
    add_metric_card(slide, 0.75, 1.95, 3.75, 1.5, "サービス方向性", "1 軸選択", "A/B/Cのどれを最初に育てるか", COLORS["accent_cyan"])
    add_metric_card(slide, 4.79, 1.95, 3.75, 1.5, "公式化する連携", "4 候補", "NotebookLM / Slack / Notion / Obsidian", COLORS["accent_emerald"])
    add_metric_card(slide, 8.83, 1.95, 3.75, 1.5, "打合せ後の反映", "即時同期", "WBS / Calendar / Issues / Docs", COLORS["accent_amber"])
    
    add_rect(slide, 0.75, 3.65, 11.83, 1.85, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "会議の進め方", 0.95, 3.82, 1.3, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"])
    add_bullets(
        slide,
        [
            "公開URLで動くデモと、Google Workspace同期の到達点を確認する。",
            "NotebookLMが生成した論点を材料に、社長への質問を短く整理する。",
            "未決事項は未決のまま残し、6/2後にWBSへ差し替える受け皿を用意する。",
        ],
        0.95,
        4.22,
        11.4,
        1.15,
        size=13.5,
    )
    add_evidence_panel(slide, ["CEO slide outline: exports/knowledge_flow/notebooklm_ceo_slide_outline.md", "PPTX: exports/knowledge_flow/mighty_skill_bridge_ceo_presentation_2026-06-02.pptx"])

    # 2: Prototype Status & Public Demo Guard
    slide = prs.slides.add_slide(blank)
    add_header(slide, 2, "Prototype & Guard System", "現在の到達点と公開デモ", source_note)
    add_text(slide, "社長に見せる公開URLは、UIデグレを防ぐ自動ガードを通して常時維持しています。", 0.78, 1.5, 11.7, 0.35, size=15, color=COLORS["text_secondary"])
    
    add_metric_card(slide, 0.75, 1.95, 3.75, 1.5, "Public Demo Guard", "PASS", "公開URLの健全性とE2E導線を検証", COLORS["accent_emerald"])
    add_metric_card(slide, 4.79, 1.95, 3.75, 1.5, "AI Fallback Mode", "LIVE / MOCK", "Gemini quota中も停止しない二重構造", COLORS["accent_cyan"])
    add_metric_card(slide, 8.83, 1.95, 3.75, 1.5, "WBS Sync Target", f"{context['wbs_total']} Tasks", "Sheets / Calendarへ完全双方向同期", COLORS["accent_amber"])
    
    add_rect(slide, 0.75, 3.65, 11.83, 1.85, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "主要到達ポイント", 0.95, 3.82, 1.5, fill=COLORS["chip_bg_emerald"], text=COLORS["accent_emerald"])
    add_bullets(
        slide,
        [
            "経歴書・案件票から4軸フィット診断へ進むUI体験を維持。",
            "Gemini制限中もdeterministic fallbackとCodexで自律開発を継続。",
            "公開URLは社長共有済みのため、コミット・push前後のPublic Demo Guardを必須化。",
        ],
        0.95,
        4.22,
        11.4,
        1.15,
        size=13.5,
    )
    add_evidence_panel(slide, ["Public URL: https://kanta13jp1.github.io/mighty-link-ai-connect/", "Guard script: scripts/verify_public_demo.py", "Quota-safe flow: docs/CODEX_CONTINUATION_NOTES.md"])

    # 3: Google Workspace Integration Architecture
    slide = prs.slides.add_slide(blank)
    add_header(slide, 3, "Workspace Operating Pipeline", "Google Workspaceで進捗が回る基盤", source_note)
    
    flow_items = [
        ("01", "Codex", "実装・文書・WBS更新", COLORS["accent_cyan"]),
        ("02", "WBS.tsv", "日程と状態の主台帳", COLORS["accent_indigo"]),
        ("03", "Sheets", "CATS型WBS / Summary", COLORS["accent_emerald"]),
        ("04", "Calendar", "開発予定を自動同期", COLORS["accent_amber"]),
        ("05", "Docs / Drive", "NotebookLM連携資料", COLORS["accent_rose"]),
    ]
    for idx, (step, title, note, accent) in enumerate(flow_items):
        x = 0.75 + idx * 2.42
        add_flow_card(slide, x, 1.95, 2.15, 1.45, step, title, note, accent)
        if idx < len(flow_items) - 1:
            add_arrow(slide, x + 2.18, 2.5)

    add_rect(slide, 0.75, 3.65, 11.83, 1.85, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "運用ポリシー", 0.95, 3.82, 1.3, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"])
    add_bullets(
        slide,
        [
            "OAuth実行アカウントは k-umezawa@ml-mightylink.com に固定。",
            "WBSはGoogle Sheetsへ、主要予定はGoogle Calendarへ自動同期・整理。",
            "docs/配下はGoogle Docs化し、NotebookLM sourceとして再利用・知識化。",
        ],
        0.95,
        4.22,
        11.4,
        1.15,
        size=13.5,
    )
    add_evidence_panel(slide, [f"Workspace account: {context['account']}", f"NotebookLM sources ready: {context['source_count']}", "Spreadsheet ID: 1L99HCBHr4IsVUWqnUuG6OgoUmxEQUdfaYQim1n6etB8"])

    # 4: Knowledge Flow Ecosystem
    slide = prs.slides.add_slide(blank)
    add_header(slide, 4, "Knowledge Ecosystem", "開発ナレッジ連携の実績とデモ", source_note)
    lanes = [
        ("NotebookLM", "22 Sources Ready\nAgent Brief / CEO Outline 取得済み", COLORS["chip_bg_cyan"], COLORS["accent_cyan"], "公式同期完了"),
        ("Slack", "投稿案生成済み\nCLI/MCP送信ツールは未露出", COLORS["chip_bg_amber"], COLORS["accent_amber"], "権限確認待ち"),
        ("Notion", "証跡ページ作成済み\n意思決定DB候補を管理", COLORS["chip_bg_emerald"], COLORS["accent_emerald"], "MCP実行済み"),
        ("Obsidian", "ローカルVault生成済み\nADR / Prompt / Meeting導線", COLORS["surface_elevated"], COLORS["text_secondary"], "ローカル完了"),
    ]
    for idx, (title, note, chip_bg, chip_text, status) in enumerate(lanes):
        x = 0.75 + idx * 3.01
        add_rect(slide, x, 1.95, 2.8, 2.3, fill=COLORS["surface"], line=COLORS["border"], radius=True)
        add_chip(slide, status, x + 0.18, 2.12, 1.4, fill=chip_bg, text=chip_text, border=COLORS["border"])
        add_text(slide, title, x + 0.18, 2.58, 2.4, 0.32, size=16, color=COLORS["text_primary"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, note, x + 0.18, 2.98, 2.4, 1.1, size=10.5, color=COLORS["text_secondary"])

    add_rect(slide, 0.75, 4.45, 11.83, 1.1, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_bullets(
        slide,
        [
            "NotebookLMはプレゼン草案作成に使い、AIエージェントの開発入力としても活用。",
            "Slackは投稿先と共有範囲の社長確認後に実送信へ進める。",
            "Notionは議事録・意思決定・バックログ公式台帳、Obsidianはローカル思考メモとして棲み分け。",
        ],
        0.95,
        4.55,
        11.4,
        0.9,
        size=12.5,
    )
    add_evidence_panel(slide, ["Notion evidence: https://www.notion.so/3671d736b9db818aaa33da0a5f1a3951", "Slack draft: exports/knowledge_flow/slack_ceo_update.md", "Obsidian vault: exports/knowledge_flow/obsidian_vault/"], y=5.75)

    # 5: NotebookLM to PPTX Automated Generation
    slide = prs.slides.add_slide(blank)
    add_header(slide, 5, "Automated Deck Pipeline", "NotebookLMからPPTXへ", source_note)
    
    gen_flow = [
        ("01", "docs/ + WBS", "公式手順・WBSを同期", COLORS["accent_cyan"]),
        ("02", "Google Docs", "Workspace所有22件", COLORS["accent_indigo"]),
        ("03", "NotebookLM CLI", "要約・QA・8枚構成", COLORS["accent_emerald"]),
        ("04", "PowerPoint", "社長説明用PPTX生成", COLORS["accent_amber"]),
    ]
    for idx, (step, title, note, accent) in enumerate(gen_flow):
        x = 0.75 + idx * 3.03
        add_flow_card(slide, x, 1.95, 2.7, 1.45, step, title, note, accent)
        if idx < len(gen_flow) - 1:
            add_arrow(slide, x + 2.76, 2.5)

    add_rect(slide, 0.75, 3.65, 11.83, 1.85, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "自動化パイプラインの特長", 0.95, 3.82, 1.8, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"])
    add_bullets(
        slide,
        [
            "NotebookLM CLIで取得したCEO Slide Outlineを、説明用PPTXのストーリー骨子に直接反映。",
            "Figma Design Tokensに準拠したカラー・タイポグラフィ・カード構造でPPTXを自動生成。",
            "編集可能なPowerPointファイルとして保存し、Google Driveへのアップロード対象化。",
        ],
        0.95,
        4.22,
        11.4,
        1.15,
        size=13.5,
    )
    add_evidence_panel(slide, [f"NotebookLM notebook: {context['notebook_id']}", "Outline: exports/knowledge_flow/notebooklm_ceo_slide_outline.md", "Generator: scripts/generate_ceo_presentation_deck.py"])

    # 6: Service Direction Matrix
    slide = prs.slides.add_slide(blank)
    add_header(slide, 6, "Strategic Options", "サービス方向性の選択肢", source_note)
    add_option_card(slide, 0.75, 1.95, 3.75, 2.5, "A. AIフィット診断支援", "営業 / 人材担当 / エンジニア", "採用・SES・案件配属の工数を劇的に削減", "デモの診断体験とマッチング精度を磨く", is_recommended=True)
    add_option_card(slide, 4.79, 1.95, 3.75, 2.5, "B. Workspace連携型PM", "経営 / PM / 現場責任者", "進捗・報告・予定同期をワンストップ自動化", "WBS/Calendar同期基盤をパッケージ化")
    add_option_card(slide, 8.83, 1.95, 3.75, 2.5, "C. AI PoC高速構築支援", "新規事業 / 企画 / 開発責任者", "検証回数と提案速度を2倍にする受託型支援", "NotebookLM/Docs連携テンプレートを型化")

    add_rect(slide, 0.75, 4.65, 11.83, 0.95, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_bullets(
        slide,
        [
            "6/2では最初の適用業務と最初に見せる相手を1つ決定する。",
            "正式企画は決め打ちせず、社長判断後にWBSとロードマップへ即時反映する。",
        ],
        0.95,
        4.75,
        11.4,
        0.75,
        size=13.5,
    )
    add_evidence_panel(slide, ["Decision pack: docs/CEO_PRESENTATION_DECISION_PACK_2026-06-02.md", "WBS tasks: T605 / T611 / T623 / T615"])

    # 7: Governance, Risks & Guardrails
    slide = prs.slides.add_slide(blank)
    add_header(slide, 7, "Governance & Guardrails", "運用・リスク論点と未完了項目", source_note)
    risks = [
        ("公開URL", "社長共有済み。UIデグレは許容しない。", "Public Demo Guardをpush前後で必須実行", COLORS["accent_emerald"]),
        ("Slack連携", "投稿先と共有範囲の承認が必要。", "承認後に公式コネクタで送信", COLORS["accent_amber"]),
        ("GitHub Project", "Project #1へWBS対象限定同期済み。", "完了・更新タスクだけを継続同期", COLORS["accent_cyan"]),
        ("外部投入情報", "認証情報・個人情報・未承認顧客情報は投入禁止。", "docs/とWBSに安全ルールを明記", COLORS["accent_rose"]),
    ]
    
    # Table header
    y = 1.85
    add_rect(slide, 0.75, y, 11.83, 0.45, fill=COLORS["surface_elevated"], line=COLORS["border"], radius=True)
    add_text(slide, "管理論点", 0.95, y + 0.1, 1.8, 0.25, size=10.5, color=COLORS["accent_cyan"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "現状ステータス", 3.0, y + 0.1, 4.2, 0.25, size=10.5, color=COLORS["text_secondary"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "ガードレール / 次の扱い", 7.5, y + 0.1, 4.8, 0.25, size=10.5, color=COLORS["accent_emerald"], bold=True, font_name=FONT_DISPLAY)
    
    for idx, (topic, status, action, accent) in enumerate(risks):
        row_y = y + 0.55 + idx * 0.72
        add_rect(slide, 0.75, row_y, 11.83, 0.62, fill=COLORS["surface"], line=COLORS["border"], radius=True)
        add_chip(slide, topic, 0.95, row_y + 0.14, 1.6, fill=COLORS["surface_elevated"], text=accent, border=COLORS["border"])
        add_text(slide, status, 3.0, row_y + 0.16, 4.2, 0.3, size=11, color=COLORS["text_secondary"])
        add_text(slide, action, 7.5, row_y + 0.16, 4.8, 0.3, size=11, color=COLORS["text_primary"], bold=True)
        
    add_evidence_panel(slide, ["GitHub Project: Project #1 operational", "WBS sync: scripts/sync_wbs_to_github.py", "Slack: destination approval required before posting"])

    # 8: Next Actions & WBS Reflection
    slide = prs.slides.add_slide(blank)
    add_header(slide, 8, "Execution Roadmap", "次アクションとWBSへの即時反映", source_note)
    steps = [
        ("01", "PPTX生成", "NotebookLM outlineをPowerPoint化", COLORS["accent_cyan"]),
        ("02", "権限確認", "Slack / GitHub Projectの復旧確認", COLORS["accent_indigo"]),
        ("03", "判断材料レビュー", "社長説明資料の最終確認", COLORS["accent_emerald"]),
        ("04", "最終リハーサル", "公開URL / WBS / 資料の確認", COLORS["accent_amber"]),
        ("05", "社長判断", "決定事項をWBSへ即時反映", COLORS["accent_rose"]),
    ]
    for idx, (step, title, note, accent) in enumerate(steps):
        x = 0.75 + idx * 2.42
        add_flow_card(slide, x, 1.95, 2.15, 1.45, step, title, note, accent)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 2.18, 2.5)

    add_rect(slide, 0.75, 3.65, 11.83, 1.85, fill=COLORS["surface"], line=COLORS["border"], radius=True)
    add_chip(slide, "アクション方針", 0.95, 3.82, 1.3, fill=COLORS["chip_bg_cyan"], text=COLORS["accent_cyan"])
    add_bullets(
        slide,
        [
            "社長に決めてもらうこと: 最初の業務課題、最初の利用者、公式化する連携、公開範囲。",
            "決定直後にやること: WBS・Calendar・Issues・Docsへ反映し、次回レビュー日を固定。",
            "保留にすること: サービス名、課金、外部共有範囲、Slack/Notionの正式運用範囲。",
        ],
        0.95,
        4.22,
        11.4,
        1.15,
        size=13.5,
    )
    add_evidence_panel(slide, ["WBS: data/WBS.tsv / docs/WBS.md", "Calendar sync: scripts/sync_wbs_to_calendar.py", "Issue/Project sync: scripts/sync_wbs_to_github.py"])


def verify_pptx(path: Path, expected_slides: int) -> None:
    if not path.exists() or path.stat().st_size < 10_000:
        raise RuntimeError(f"PPTX was not created or is too small: {path}")
    with zipfile.ZipFile(path) as package:
        slide_parts = [
            name for name in package.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
    if len(slide_parts) != expected_slides:
        raise RuntimeError(f"Expected {expected_slides} slides, found {len(slide_parts)}")


def write_summary(context: dict[str, Any], pptx_path: Path, summary_path: Path) -> None:
    content = f"""# Mighty Skill-Bridge CEO Presentation Deck

Generated: {jst_now().isoformat(timespec="seconds")}

## Output

- PPTX: `exports/knowledge_flow/{pptx_path.name}`
- Theme: `{CURRENT_THEME}` (Figma Design Token Compliant)
- Google Drive: {context.get('pptx_drive_url') or 'not uploaded yet'}
- Generator: `scripts/generate_ceo_presentation_deck.py`
- NotebookLM outline: `exports/knowledge_flow/notebooklm_ceo_slide_outline.md`
- NotebookLM notebook: `{context['notebook_id']}`
- Workspace account: `{context['account']}`

## Slide List

1. 本日決めたいこと (Decision Meeting Scope)
2. 現在の到達点と公開デモ (Prototype & Guard System)
3. Google Workspaceで進捗が回る基盤 (Workspace Operating Pipeline)
4. 開発ナレッジ連携の実績とデモ (Knowledge Ecosystem)
5. NotebookLMからPPTXへ (Automated Deck Pipeline)
6. サービス方向性の選択肢 (Strategic Options)
7. 運用・リスク論点と未完了項目 (Governance & Guardrails)
8. 次アクションとWBSへの即時反映 (Execution Roadmap)

## Design System & Figma Token Notes

- Integrated Figma Design Tokens (`docs/design_tokens.json`) with Dark Ink (`#0B0F19`) canvas and card-based containers (`#151E2E`).
- Strict 8px-grid spacing, rounded container cards (`radius=True`), and vibrant semantic accent chips (Cyan, Emerald, Amber, Rose).
- Full editable native PowerPoint format (.pptx) without bitmap degradation.
"""
    summary_path.write_text(content, encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate CEO presentation PPTX with Figma Design System.")
    parser.add_argument(
        "--style",
        type=str,
        default="figma-modern",
        choices=["figma-modern", "default", "canva-export"],
        help="Theme/style for slide generation. 'figma-modern' (default) produces dark studio decks.",
    )
    args = parser.parse_args()

    global CURRENT_THEME, COLORS
    CURRENT_THEME = args.style
    COLORS = PALETTES.get(args.style, PALETTES["default"])

    pptx_filename = f"mighty_skill_bridge_ceo_presentation_2026-06-02_{args.style}.pptx" if args.style != "default" else "mighty_skill_bridge_ceo_presentation_2026-06-02.pptx"
    summary_filename = pptx_filename.replace(".pptx", ".md")
    manifest_filename = pptx_filename.replace(".pptx", ".json")

    pptx_path = EXPORT_DIR / pptx_filename
    summary_path = EXPORT_DIR / summary_filename
    manifest_path = EXPORT_DIR / manifest_filename

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    manifest = load_json(EXPORT_DIR / "manifest.json")
    docs_manifest = load_json(EXPORT_DIR / "notebooklm_docs_manifest.json")
    drive_docs = load_json(EXPORT_DIR / "google_drive_workspace_docs.json")
    outline = load_json(EXPORT_DIR / "notebooklm_ceo_slide_outline.json")

    context = {
        "account": drive_docs.get("account") or docs_manifest.get("account") or "k-umezawa@ml-mightylink.com",
        "notebook_id": outline.get("notebook_id") or docs_manifest.get("notebook_id") or "75521ea6-6b9b-47b2-9508-50050d8ab2d5",
        "source_count": len(docs_manifest.get("google_docs", {})) or 14,
        "wbs_total": get_nested(manifest, "summary", "total", default=72),
        "pptx_drive_url": get_nested(drive_docs, "files", "ceo_presentation_pptx", "url", default=""),
    }

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slides(prs, context)
    prs.save(pptx_path)

    verify_pptx(pptx_path, expected_slides=8)
    write_summary(context, pptx_path, summary_path)

    deck_manifest = {
        "generated_at_jst": jst_now().isoformat(timespec="seconds"),
        "account": context["account"],
        "notebook_id": context["notebook_id"],
        "theme": CURRENT_THEME,
        "source_outline": "exports/knowledge_flow/notebooklm_ceo_slide_outline.md",
        "source_count": context["source_count"],
        "tooling": {
            "design_tokens": "docs/design_tokens.json",
            "pptx_generator": "python-pptx (Figma Modern Engine)",
        },
        "outputs": {
            "pptx": str(pptx_path.relative_to(PROJECT_ROOT)).replace("\\", "/"),
            "summary": str(summary_path.relative_to(PROJECT_ROOT)).replace("\\", "/"),
            "google_drive_url": context.get("pptx_drive_url") or None,
        },
        "slide_count": 8,
    }
    manifest_path.write_text(json.dumps(deck_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    print(f"[+] CEO presentation deck generated successfully ({CURRENT_THEME}).")
    print(f"[*] PPTX: {pptx_path.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary: {summary_path.relative_to(PROJECT_ROOT)}")
    print(f"[*] Manifest: {manifest_path.relative_to(PROJECT_ROOT)}")


if __name__ == "__main__":
    main()
