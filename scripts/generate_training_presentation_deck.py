#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate Premium Canva-Inspired Presentation Deck for 2026-08-26 Internal AI Training.

Theme: 'Google Antigravity 1-Hour Live Demo & Hands-on'
Design Style: Canva Modern Studio (Mint & Deep Ink, Glass Card Containers, Accent Badges)
Slide Aspect Ratio: 16:9 Wide (1920x1080 / 13.333x7.5 in)

Agenda Structure (60-Minute Live Demo):
1. Title & Hero Cover: AIエージェントでWeb制作・アプリ開発を1時間で体験する (Antigravity Live Demo)
2. Agenda & 60-Minute Timeline: 1時間のタイムテーブル（イントロ10分 / ライブデモ25分 / 参加者ハンズオン15分 / Q&A 10分）
3. Why Antigravity: なぜ今、Google Antigravityなのか（Artifacts × Visual Feedback × Browser Agent）
4. Live Demo Flow (25 min): ゼロからWebアプリを完成させる4ステップ
5. Visual Feedback Feature: 画面上の直接ピン留め＆コメント指示の威力
6. 9 AI Agents Comparison: 他のAIエージェント（Claude Code / Codex等）との棲み分け
7. Safe Usage & Guardrails: 社内利用ルール・機密情報保護
8. Summary & Next Steps: 本日のゴール、社内申請・テンプレート活用
"""

from __future__ import annotations

import datetime as dt
import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports" / "training_deck"
PPTX_FILE = EXPORT_DIR / "internal_ai_training_demo_2026-08-26.pptx"
SUMMARY_FILE = EXPORT_DIR / "internal_ai_training_demo_2026-08-26.md"
CANVA_GUIDE_FILE = EXPORT_DIR / "CANVA_IMPORT_GUIDE.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Canva Clean Studio Theme Palette (Mint & Slate)
COLORS = {
    "bg_dark": RGBColor(11, 15, 25),          # Deep Canvas #0B0F19
    "card_dark": RGBColor(21, 30, 46),        # Surface Slate #151E2E
    "card_elevated": RGBColor(30, 41, 59),    # Highlight Container #1E293B
    "border_subtle": RGBColor(42, 54, 79),    # Divider line #2A364F
    
    # Canva Mint & Gradient Accents
    "mint_bright": RGBColor(45, 212, 191),    # Teal/Mint 400 #2DD4BF
    "mint_bg": RGBColor(13, 53, 50),          # Mint Deep Container
    "sky_cyan": RGBColor(56, 189, 248),       # Sky 400 #38BDF8
    "purple_accent": RGBColor(167, 139, 250), # Purple 400 #A78BFA
    "amber_accent": RGBColor(251, 191, 36),   # Amber 400 #FBBF24
    "emerald_accent": RGBColor(52, 211, 153), # Emerald 400 #34D399
    "rose_accent": RGBColor(251, 113, 133),   # Rose 400 #FB7185
    
    # Text
    "text_title": RGBColor(248, 250, 252),    # #F8FAFC
    "text_body": RGBColor(203, 213, 225),     # #CBD5E1
    "text_muted": RGBColor(148, 163, 184),    # #94A3B8
    "text_subtle": RGBColor(100, 116, 139),   # #64748B
}

FONT_DISPLAY = "Plus Jakarta Sans"
FONT_BODY = "Yu Gothic"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name or FONT_BODY
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color or COLORS["text_title"]
    return box


def add_bullets(
    slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 13.5,
    color: RGBColor | None = None,
    font_name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"✔  {bullet}"
        p.level = 0
        p.space_after = Pt(7)
        run = p.runs[0]
        run.font.name = font_name or FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["text_body"]
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_chip(
    slide,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    text: RGBColor,
    border: RGBColor | None = None,
):
    add_rect(slide, x, y, w, 0.34, fill=fill, line=border or fill, radius=True)
    add_text(
        slide,
        label,
        x,
        y + 0.04,
        w,
        0.24,
        size=9.5,
        color=text,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        font_name=FONT_DISPLAY,
    )


def add_header(slide, slide_no: int, eyebrow: str, title: str):
    # Full dark canvas
    add_rect(slide, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], line=COLORS["bg_dark"])
    
    # Canva Modern Dual Pill Bar
    add_rect(slide, 0.75, 0.35, 0.6, 0.05, fill=COLORS["mint_bright"], line=COLORS["mint_bright"])
    add_rect(slide, 1.4, 0.35, 0.35, 0.05, fill=COLORS["sky_cyan"], line=COLORS["sky_cyan"])
    
    # Eyebrow
    add_text(
        slide,
        f"{slide_no:02d}  /  {eyebrow.upper()}",
        0.75,
        0.48,
        8.0,
        0.26,
        size=9,
        color=COLORS["mint_bright"],
        bold=True,
        font_name=FONT_MONO,
    )
    # Title
    add_text(
        slide,
        title,
        0.73,
        0.76,
        11.8,
        0.65,
        size=26,
        color=COLORS["text_title"],
        bold=True,
        font_name=FONT_DISPLAY,
    )
    
    # Footer
    add_text(
        slide,
        f"2026-08-26 社内AI研修 · Antigravity 1時間ライブデモ  ·  Slide {slide_no:02d} / 08",
        0.75,
        7.05,
        7.0,
        0.28,
        size=8.5,
        color=COLORS["text_subtle"],
        font_name=FONT_MONO,
    )
    add_text(
        slide,
        "Canva Compatible Deck  ·  Mighty Link AI Connect",
        8.0,
        7.05,
        4.58,
        0.28,
        size=8.5,
        color=COLORS["text_subtle"],
        align=PP_ALIGN.RIGHT,
        font_name=FONT_DISPLAY,
    )


def build_slides(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]

    # --- Slide 1: Cover & Hero ---
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], line=COLORS["bg_dark"])
    
    # Main Hero Frame (Canva Mint/Dark Card)
    add_rect(slide, 0.75, 0.85, 11.83, 5.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    
    add_chip(slide, "2026-08-26 オンライン開催 · 60分ライブデモ", 1.2, 1.25, 4.2, fill=COLORS["mint_bg"], text=COLORS["mint_bright"], border=COLORS["mint_bright"])
    
    add_text(slide, "AIエージェントでWeb制作・アプリ開発を", 1.15, 1.75, 11.0, 0.65, size=30, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "1時間で体験する Antigravity ライブデモ", 1.15, 2.4, 11.0, 0.65, size=30, color=COLORS["mint_bright"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "〜 非エンジニアと開発者が共に創る「Artifacts × 画面直接指示」の次世代開発体験 〜", 1.18, 3.15, 11.0, 0.35, size=15, color=COLORS["text_body"])
    
    # 3 Key Highlights Cards
    highlights = [
        ("01. 自然言語で即時モック化", "「こんな機能が欲しい」と入力するだけで、AIが設計・コード・画面を即時生成", COLORS["mint_bright"]),
        ("02. Visual Feedback修正", "画面の気になる部分を直接クリックしてピン留め。言葉足らずの指示ミスを撲滅", COLORS["sky_cyan"]),
        ("03. 1時間完結ライブ作成", "アイデア出しからWeb公開・PowerPoint資料化までの全工程を完全生実演", COLORS["purple_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(highlights):
        x = 1.18 + idx * 3.7
        add_rect(slide, x, 3.75, 3.45, 2.5, fill=COLORS["card_elevated"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.2, 4.0, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.2, 4.2, 3.05, 0.32, size=15, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.2, 4.65, 3.05, 1.3, size=11.5, color=COLORS["text_muted"])

    # --- Slide 2: 60-Minute Time Table ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 2, "Session Schedule", "本日のタイムテーブル（60分のアジェンダ）")
    add_text(slide, "理論よりも「実際の操作画面と挙動」を体感していただく実践中心のプログラムです。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    timeline = [
        ("10 min", "イントロダクション", "AIエージェントの現在地と「なぜAntigravityなのか」の背景解説", COLORS["mint_bright"]),
        ("25 min", "ゼロからWeb制作 ライブ実演", "画面共有でAntigravityを操作。要件定義からデプロイまでをノーカット実演", COLORS["sky_cyan"]),
        ("15 min", "参加者プチ・ハンズオン", "各自の端末でプロンプト入力 & Visual Feedbackピン留め指示を体験", COLORS["purple_accent"]),
        ("10 min", "質疑応答 & 社内活用案内", "社内アカウント申請手順・共有プロンプト集・今後の展開案内", COLORS["emerald_accent"]),
    ]
    for idx, (time_tag, title, desc, accent) in enumerate(timeline):
        x = 0.75 + idx * 3.01
        add_rect(slide, x, 1.95, 2.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, time_tag, x + 0.2, 2.2, 1.2, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, x + 0.2, 2.75, 2.4, 0.65, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.2, 3.55, 2.4, 2.8, size=12, color=COLORS["text_body"])

    # --- Slide 3: Why Antigravity ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 3, "Core Advantages", "なぜ今、Google Antigravityなのか？")
    add_text(slide, "従来のチャット型AI（ChatGPT/Claude）と決定的に異なる「4大パラダイムシフト」", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    reasons = [
        ("Artifacts (動的成果物)", "チャット欄に長いコードを流すのではなく、画面右側に独立したドキュメント・計画・プレビューを即時生成。いつでも修正・参照が可能。", COLORS["mint_bright"]),
        ("Visual Feedback", "「ここのボタンを青にして」「この余白を狭くして」といった指示を、画面キャプチャ上のクリック位置ピンと短いコメントだけで正確に伝達。", COLORS["sky_cyan"]),
        ("Autonomous Browser", "Headless Chromeを自動起動し、作成したWebサイトのボタンクリック・E2Eテスト・レスポンシブ崩れをAI自身が自律検証。", COLORS["purple_accent"]),
        ("外部ツール・MCP連携", "Figma、Canva、Google Workspace、Slack等の外部SaaSと双方向でデータをやり取りし、社内業務フローへ即座に統合。", COLORS["amber_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(reasons):
        row = idx // 2
        col = idx % 2
        x = 0.75 + col * 6.02
        y = 1.95 + row * 2.4
        add_rect(slide, x, y, 5.8, 2.2, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, y + 0.22, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, y + 0.38, 5.3, 0.35, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.25, y + 0.85, 5.3, 1.15, size=12, color=COLORS["text_body"])

    # --- Slide 4: 25-Minute Live Demo Flow ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 4, "Live Demonstration", "25分ライブデモの実演ステップ（実画面共有）")
    add_text(slide, "「社内向けAI比較ポータルサイト」をゼロから作成・公開する一連のフローを完全実演します。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    demo_steps = [
        ("STEP 01", "要件プロンプト投入", "「9大AIエージェント比較サイトを作って」と自然言語で指示", COLORS["mint_bright"]),
        ("STEP 02", "計画 & 自動コーディング", "Implementation Planを策定しHTML/CSS/JSを一括自動生成", COLORS["sky_cyan"]),
        ("STEP 03", "画面直接ピン留め修正", "Artifactプレビューを見ながらVisual Feedbackで直感的に微調整", COLORS["purple_accent"]),
        ("STEP 04", "即時デプロイ & 導通確認", "GitHub Pagesへ自動公開しHTTP 200疎通確認を完遂", COLORS["emerald_accent"]),
    ]
    for idx, (step_tag, title, desc, accent) in enumerate(demo_steps):
        x = 0.75 + idx * 3.01
        add_rect(slide, x, 1.95, 2.8, 3.2, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, step_tag, x + 0.2, 2.15, 1.2, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, x + 0.2, 2.65, 2.4, 0.55, size=15, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.2, 3.3, 2.4, 1.6, size=11.5, color=COLORS["text_body"])

    # Bottom Live URL Box
    add_rect(slide, 0.75, 5.35, 11.83, 1.45, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    add_chip(slide, "🌐 実演で公開される完成サイト", 0.95, 5.52, 2.8, fill=COLORS["mint_bg"], text=COLORS["mint_bright"], border=COLORS["mint_bright"])
    add_text(slide, "URL: https://kanta13jp1.github.io/mighty-link-antigravity-live-demo/ （全9製品比較 & 公式アイコン & 動画デモシアター搭載）", 0.95, 6.0, 11.4, 0.5, size=12, color=COLORS["text_body"])

    # --- Slide 5: Visual Feedback Deep Dive ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 5, "Visual Feedback", "「言葉足らず」を解決する Visual Feedback の仕組み")
    add_text(slide, "非エンジニアでも迷わない！画面をクリックして直感的に修正指示を出す3ステップ", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    vf_cards = [
        ("1. 画面を開く", "AIが作成したWebサイトやUIモックアップのArtifactプレビューを表示します。", COLORS["mint_bright"]),
        ("2. 気になる場所をクリック", "修正したいボタンや文字を直接クリックすると、その場所に「📌 ピン」が刺さります。", COLORS["sky_cyan"]),
        ("3. 短い指示を入力して送信", "「ここを赤色にして」「文字サイズを大きく」と短いコメントを入力するだけでAIが的確に修正。", COLORS["purple_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(vf_cards):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 3.2, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, 2.2, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, 2.45, 3.3, 0.35, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.25, 3.0, 3.3, 1.8, size=12.5, color=COLORS["text_body"])

    # Highlight Callout
    add_rect(slide, 0.75, 5.35, 11.83, 1.45, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    add_chip(slide, "💡 従来のチャット指示との違い", 0.95, 5.52, 2.8, fill=COLORS["card_elevated"], text=COLORS["amber_accent"], border=COLORS["amber_accent"])
    add_text(slide, "「上から3つ目のカードの右下のボタン...」のような複雑な説明が一切不要になり、指示ミスや手戻りがゼロになります。", 0.95, 6.0, 11.4, 0.5, size=12.5, color=COLORS["text_body"])

    # --- Slide 6: 9 AI Agents Comparison ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 6, "AI Ecosystem", "主要AIエージェントの使い分けマップ（2026）")
    add_text(slide, "Antigravityだけでなく、用途や業務内容に応じて最適なツールを組み合わせます。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    matrix_rows = [
        ("Google Antigravity", "Web制作・画面直接指示・プロトタイピング", "非エンジニア・デザイナー・フロントエンド", "Artifacts / Visual Feedback / ブラウザ自動検証"),
        ("OpenAI Codex", "既存コードベース理解・バックエンド実装・バグ修正", "ソフトウェアエンジニア・SRE", "AGENTS.md / MCP / クラウドサンドボックス"),
        ("Anthropic Claude Code", "マルチファイル開発・Git操作・ターミナル自動化", "フルスタックエンジニア・テックリード", "CLAUDE.md / subagents / 自動メモリ"),
        ("Claude Cowork", "ドキュメント作成・表計算・リサーチ・定期レポート", "バックオフィス・企画・マーケター", "Projects / Connectors / scheduled tasks"),
    ]
    y = 1.95
    add_rect(slide, 0.75, y, 11.83, 0.45, fill=COLORS["card_elevated"], line=COLORS["border_subtle"], radius=True)
    add_text(slide, "エージェント名", 0.95, y + 0.1, 2.2, 0.25, size=10.5, color=COLORS["mint_bright"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "得意な領域・ユースケース", 3.3, y + 0.1, 3.8, 0.25, size=10.5, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "おすすめの利用者", 7.3, y + 0.1, 2.2, 0.25, size=10.5, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "主な機能・武器", 9.7, y + 0.1, 2.6, 0.25, size=10.5, color=COLORS["text_muted"], bold=True, font_name=FONT_DISPLAY)

    for idx, (name, domain, user, feat) in enumerate(matrix_rows):
        row_y = y + 0.55 + idx * 1.15
        add_rect(slide, 0.75, row_y, 11.83, 1.02, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_text(slide, name, 0.95, row_y + 0.18, 2.2, 0.3, size=13, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, domain, 3.3, row_y + 0.18, 3.8, 0.65, size=11, color=COLORS["text_body"])
        add_text(slide, user, 7.3, row_y + 0.18, 2.2, 0.65, size=11, color=COLORS["mint_bright"])
        add_text(slide, feat, 9.7, row_y + 0.18, 2.6, 0.65, size=10.5, color=COLORS["text_muted"])

    # --- Slide 7: Safe Usage & Security Guardrails ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 7, "Security & Rules", "社内利用における安全ガイドラインと禁止事項")
    add_text(slide, "全社で安心してAIを活用するために徹底すべき4つのセキュリティ原則", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    rules = [
        ("RULE 1: 機密情報・個人情報の投入禁止", "顧客の実名、パスワード、APIキー、未公開の契約書テキストはプロンプトに投入厳禁。テスト時はダミーデータを使用する。", COLORS["rose_accent"]),
        ("RULE 2: 生成物のファクトチェック必須", "AIの回答や生成コードにはハルシネーション（誤情報）が含まれうるため、必ず担当者が動作・内容を確認する。", COLORS["amber_accent"]),
        ("RULE 3: 外部公開時の著作権・OSSライセンス確認", "生成した画像・文章・ソースコードを社外へ公開・納品する場合は、ライセンスおよび商用利用規約を確認する。", COLORS["sky_cyan"]),
        ("RULE 4: アカウント・クォータの適正利用", "各エージェントに設定された利用上限（クォータ）を意識し、大規模な自動処理を行う際は事前に管理部へ相談する。", COLORS["emerald_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(rules):
        row_y = 1.95 + idx * 1.18
        add_rect(slide, 0.75, row_y, 11.83, 1.05, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, f"SECURITY 0{idx+1}", 0.95, row_y + 0.15, 1.4, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, 2.55, row_y + 0.15, 9.6, 0.3, size=13.5, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, 2.55, row_y + 0.48, 9.6, 0.45, size=11, color=COLORS["text_body"])

    # --- Slide 8: Next Steps & Q&A ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 8, "Next Actions", "本日のまとめ & 研修後の実践ステップ")
    add_text(slide, "研修終了後、各自の業務で小さなAI活用（Quick Win）をスタートしましょう！", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    steps = [
        ("STEP 1: アカウント申請", "社内ポータルからAntigravity / Claudeの利用申請を行う（即日承認）", COLORS["mint_bright"]),
        ("STEP 2: 社内プロンプト集の活用", "共有リポジトリに登録されている「業務別プロンプト・Skills定義」を試す", COLORS["sky_cyan"]),
        ("STEP 3: 成果共有と勉強会", "月次開催の「AI活用LT会」で自作プロンプトや自動化の成果を発表", COLORS["purple_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(steps):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 3.2, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, 2.2, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, 2.45, 3.3, 0.35, size=15, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.25, 3.0, 3.3, 1.8, size=12, color=COLORS["text_body"])

    # Contact & Q&A Box
    add_rect(slide, 0.75, 5.35, 11.83, 1.45, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    add_chip(slide, "💬 質疑応答 & 社内サポート窓口", 0.95, 5.52, 2.8, fill=COLORS["mint_bg"], text=COLORS["mint_bright"], border=COLORS["mint_bright"])
    add_text(slide, "社内Slack: #ai-agent-hub  |  質問窓口: ai-support@ml-mightylink.com  |  研修アンケートへのご協力をお願いします", 0.95, 6.0, 11.4, 0.5, size=11.5, color=COLORS["text_body"])


def verify_pptx(path: Path, expected_slides: int) -> None:
    if not path.exists() or path.stat().st_size < 10_000:
        raise RuntimeError(f"PPTX was not created or is too small: {path}")
    with zipfile.ZipFile(path) as package:
        slide_parts = [
            name for name in package.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
    if len(slide_parts) != expected_slides:
        raise RuntimeError(f"Expected {expected_slides} slides, found {len(slide_parts)}")


def write_summary_and_guide() -> None:
    summary_content = f"""# 2026-08-26 社内向けAI活用研修 デモ資料 (Canva Studioスタイル)

生成日時: {jst_now().isoformat(timespec="seconds")}

## 成果物ファイル
- **PowerPoint (Canvaインポート用PPTX)**: `exports/training_deck/{PPTX_FILE.name}`
- **Canvaインポート手順書**: `exports/training_deck/CANVA_IMPORT_GUIDE.md`

## スライド構成 (全8枚 / 16:9 ワイド / 60分ライブデモ設計)
1. **表紙 & 研修概要**: AIエージェントでWeb制作・アプリ開発を1時間で体験する Antigravity ライブデモ
2. **タイムテーブル (60分)**: イントロ10分 / ライブデモ25分 / ハンズオン15分 / Q&A 10分
3. **Antigravityの強み**: Artifacts × Visual Feedback × Autonomous Browser
4. **25分ライブデモ実演フロー**: 要件入力 ➜ 計画 & 生成 ➜ ピン留め直接指示 ➜ 自動デプロイ
5. **Visual Feedback の威力**: 画面直接クリック＆ピン留め修正の3ステップ
6. **9大AIエージェント比較**: Antigravity, Codex, Claude Code, Cowork の適材適所
7. **セキュリティ & 社内ルール**: 機密情報投入禁止, ファクトチェック, クォータ管理
8. **まとめ & ネクストアクション**: アカウント申請, 社内プロンプト集, サポート窓口
"""
    SUMMARY_FILE.write_text(summary_content, encoding="utf-8")

    canva_guide_content = """# Canva へのインポート & 編集手順

生成されたPPTXファイル (`exports/training_deck/internal_ai_training_demo_2026-08-26.pptx`) は、Canvaのプレゼンテーションエンジン（16:9 ワイド）と100%互換性を持つようにレイアウト設計されています。

## Canva への取り込み手順 (3ステップ)

1. **Canva を開く**:
   - [Canva 公式サイト](https://www.canva.com/) にログインします。

2. **ファイルをアップロード**:
   - 右上の **「デザインを作成」** → **「ファイルをインポート」** を選択します。
   - または、ホーム画面に `internal_ai_training_demo_2026-08-26.pptx` を直接ドラッグ＆ドロップします。

3. **Canva 上で自由に編集・発表**:
   - テキスト、ミント＆ダークのカード図形、カラーがすべて編集可能なレイアウトとして開きます。
   - Canvaの豊富なイラスト・アイコン・アニメーション・発表者ノートを追加して、8/26の研修本番にご活用いただけます！
"""
    CANVA_GUIDE_FILE.write_text(canva_guide_content, encoding="utf-8")


def main() -> None:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    build_slides(prs)
    prs.save(PPTX_FILE)
    
    verify_pptx(PPTX_FILE, expected_slides=8)
    write_summary_and_guide()
    
    print("[+] 2026-08-26 社内向けAI研修デモ資料 (Canva Studioスタイル・全8枚) を生成しました。")
    print(f"[*] PPTX: {PPTX_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary: {SUMMARY_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Canva Guide: {CANVA_GUIDE_FILE.relative_to(PROJECT_ROOT)}")


if __name__ == "__main__":
    main()
