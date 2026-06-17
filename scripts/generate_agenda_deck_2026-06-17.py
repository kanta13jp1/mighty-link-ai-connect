#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate a Canva/Figma-style branded CEO PPTX for the 2026-06-17 Meeting using python-pptx.

Mighty Skill-Bridge brand palette (Seedance cinematic theme):
- bg: cyber black #0d0e15
- panel: deep navy #161824
- accent1: neon blue #00f0ff
- accent2: neon green #39ff14
- accent3: neon red (risk) #ff3366
- text primary: cool white #f1f5ff
- text secondary: gray white #c5cae0

Output: exports/mighty_skill_bridge_agenda_2026-06-17.pptx
"""

from __future__ import annotations

import datetime as dt
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports"
PPTX_FILE = EXPORT_DIR / "mighty_skill_bridge_agenda_2026-06-17.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Mighty Skill-Bridge brand palette (Seedance cinematic)
C = {
    "bg": RGBColor(0x0D, 0x0E, 0x15),
    "panel": RGBColor(0x16, 0x18, 0x24),
    "panel_alt": RGBColor(0x1F, 0x22, 0x33),
    "neon_blue": RGBColor(0x00, 0xF0, 0xFF),
    "neon_green": RGBColor(0x39, 0xFF, 0x14),
    "neon_red": RGBColor(0xFF, 0x33, 0x66),
    "neon_yellow": RGBColor(0xFF, 0xD7, 0x00),
    "text_primary": RGBColor(0xF1, 0xF5, 0xFF),
    "text_secondary": RGBColor(0xC5, 0xCA, 0xE0),
    "text_muted": RGBColor(0x7A, 0x83, 0x99),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_HEADER = "Yu Gothic UI"
FONT_BODY = "Yu Gothic UI"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h, *,
             font_size=14, color=None, bold=False,
             font_name=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    color = color or C["text_primary"]
    font_name = font_name or FONT_BODY
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def base_slide(prs):
    """Create a blank slide with cyber-black background and corner accent."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Cyber black background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C["bg"])
    # Top-left neon corner accent (5px line, 40% width)
    add_rect(slide, Inches(0), Inches(0), Inches(5.2), Inches(0.04), C["neon_blue"])
    # Bottom-right neon green dot accent
    add_rect(slide, Inches(13.0), Inches(7.2), Inches(0.18), Inches(0.18), C["neon_green"])
    return slide


def add_header(slide, slide_num, title, subtitle=None):
    """Top header with slide number, title, and accent line."""
    # Slide number badge
    add_text(slide, f"0{slide_num}", Inches(0.5), Inches(0.35), Inches(0.8), Inches(0.5),
             font_size=14, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Title
    add_text(slide, title, Inches(1.3), Inches(0.3), Inches(11.5), Inches(0.6),
             font_size=24, color=C["text_primary"], bold=True, font_name=FONT_HEADER,
             align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(1.3), Inches(0.85), Inches(11.5), Inches(0.4),
                 font_size=12, color=C["neon_green"], font_name=FONT_BODY,
                 align=PP_ALIGN.LEFT)
    # Header underline
    add_rect(slide, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.02), C["neon_blue"])


def add_footer(slide, label="Mighty Skill-Bridge · 2026-06-17 CEO Regular Meeting Agenda"):
    """Bottom footer with brand label."""
    add_text(slide, label, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.3),
             font_size=9, color=C["text_muted"], font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)


def add_cta_box(slide, question, x=Inches(0.5), y=Inches(6.1), w=Inches(12.3), h=Inches(0.75)):
    """Highlighted question box for 社長への質問."""
    add_rect(slide, x, y, w, h, C["panel"], C["neon_blue"], Pt(1.5))
    # Q label
    add_text(slide, "▶ 打合せ合意・相談事項", Inches(0.75), y + Emu(50000), Inches(3.5), Inches(0.3),
             font_size=10, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Question text
    add_text(slide, question, Inches(0.75), y + Inches(0.32), w - Inches(0.5), Inches(0.42),
             font_size=13, color=C["text_primary"], bold=True, font_name=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_bullets(slide, points, x, y, w, h, *,
                font_size=14, color=None, bullet_color=None):
    color = color or C["text_secondary"]
    bullet_color = bullet_color or C["neon_green"]
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for idx, point in enumerate(points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if idx > 0:
            p.space_before = Pt(8)
        # Bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "■ "
        run_bullet.font.name = FONT_MONO
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        # Body
        run_body = p.add_run()
        run_body.text = point
        run_body.font.name = FONT_BODY
        run_body.font.size = Pt(font_size)
        run_body.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, *, label=None, accent=None):
    """Decorated panel (dark navy bg with left neon stripe)."""
    accent = accent or C["neon_blue"]
    add_rect(slide, x, y, w, h, C["panel"])
    add_rect(slide, x, y, Inches(0.08), h, accent)
    if label:
        add_text(slide, label, x + Inches(0.25), y + Inches(0.15),
                 w - Inches(0.4), Inches(0.3),
                 font_size=10, color=accent, bold=True, font_name=FONT_MONO,
                 align=PP_ALIGN.LEFT)


SLIDES = [
    {
        "num": 1,
        "title": "本日のアジェンダ & 意思決定項目",
        "subtitle": "管理下デモGo / 一般公開No-Goの確認と優先開発ロードマップ合意",
        "points": [
            "1. 社内パイロット（6/3〜6/16）の運用実績・評価報告",
            "2. コスト台帳（$0枠維持）およびQ2セキュリティ・プライバシー監査結果",
            "3. 本番ドメイン（mightylink-app.com）接続完了と費用0円化の報告",
            "4. 限定デモ継続、一般公開No-Go、営業メールAIマッチング優先の合意",
        ],
        "evidence": "WBSタスク T746 / docs/PILOT_REPORT_2026-06-16.md",
        "question": "本日のゴールとして「限定デモ継続と次フェーズ優先機能の合意」でよろしいでしょうか？",
        "accent": "neon_blue",
    },
    {
        "num": 2,
        "title": "社内パイロット運用実績 (6/3～6/16)",
        "subtitle": "AIフィット診断の実稼働とユーザー評価",
        "points": [
            "人材担当・営業担当の協力を得て『Mighty Skill-Bridge』を実稼働検証",
            "サンプルデータ（経歴書5件・案件票5件）の流し込みと4次元適合率UIの正常動作を確認",
            "診断結果に基づく面談想定質問の自動生成により、担当者の準備工数を約70%削減できる見込み",
        ],
        "evidence": "実績資料: docs/PILOT_REPORT_2026-06-16.md",
        "question": "パイロット運用の実績および約70%の工数削減効果について、ご質問やフィードバックはありますか？",
        "accent": "neon_green",
    },
    {
        "num": 3,
        "title": "コスト台帳＆セキュリティ・プライバシー監査結果",
        "subtitle": "完全無料枠内での運用維持とセキュリティの担保",
        "points": [
            "Gemini API キャッシュ（TTL 1時間）の活用により、トークン消費量を90%削減しAPI費用 $0を維持",
            "個人情報保護・GDPRに準拠し、データはパイロット終了後3営業日以内に完全消去する運用を徹底",
            "Q2セキュリティ監査をクリア、Basic Authによる第三者アクセス防止およびAPI自動遮断（ガード）を実装済",
        ],
        "evidence": "docs/COST_REPORT_2026-06.md\ndocs/SECURITY_AUDIT_REPORT_2026-Q2.md\ndata/secret_rotation_inventory.tsv",
        "question": "利用同意書の運用（3営業日以内の完全消去）および現在のコストセーフガード方針で進めてよろしいでしょうか？",
        "accent": "neon_red",
    },
    {
        "num": 4,
        "title": "本番ドメイン接続完了と構成報告",
        "subtitle": "mightylink-app.com 接続完了（追加固定費 0 円）",
        "points": [
            "サービス専用ドメイン「mightylink-app.com」をお名前.comで取得し、Firebase Hostingへ接続完了",
            "SSL証明書はFirebaseが自動発行・更新するため、追加費用や更新管理の手間は一切なし",
            "当初予定の会社ドメインから専用ドメインへ変更したことで、会社メールやHPの停止リスクを100%回避",
        ],
        "evidence": "URL: https://mightylink-app.com/ (HTTPS疎通)\ndocs/PRODUCTION_DOMAIN_SETUP_GUIDE.md",
        "question": "専用ドメイン（mightylink-app.com）での本番URLおよび特商法表記での運用を開始してよろしいでしょうか？",
        "accent": "neon_blue",
    },
    {
        "num": 5,
        "title": "開発・運用体制とアカウント管理の今後",
        "subtitle": "個人アカウントから会社アカウントへの安全な移行計画",
        "points": [
            "現在はスピード優先のため、個人アカウント主体（GitHub, Firebase, Supabase）で構築・開発中",
            "Google Workspace/WBSおよびNotebookLMは会社提供アカウント（k-umezawa@ml-mightylink.com）で管理",
            "移行手順書（ランブック）を整備済、有料化・正式リリースの段階で会社アカウントへ安全に移管可能",
        ],
        "evidence": "docs/GOOGLE_WORKSPACE_MIGRATION_RUNBOOK.md\nscripts/share_resources.py",
        "question": "今後の正式リリースおよび有料化の段階で、会社アカウントへ移管・集約するスケジュールでよろしいでしょうか？",
        "accent": "neon_yellow",
    },
    {
        "num": 6,
        "title": "Go/No-Go判定と今後のロードマップ",
        "subtitle": "限定デモGo、一般公開/有償ローンチNo-Go、T817優先",
        "points": [
            "【判定】controlled_demo はGo。社長説明・限定デモ・管理下確認は継続可能",
            "【No-Go】public_paid_launch は法務、同意UI、Stripe、負荷テスト、営業メールAIマッチングMVP完了まで保留",
            "【最優先】毎日約1,000通の営業メールから案件要件を抽出し、Supabase DBとAIマッチングへつなぐT817を優先",
        ],
        "evidence": "WBSタスク: T746 / T817 / T817_1 / T804\ndata/release_go_no_go_criteria.tsv",
        "question": "限定デモは継続し、一般公開はNo-Goのまま、営業メールAIマッチングを最優先に進めてよろしいでしょうか？",
        "accent": "neon_green",
    },
]


def render_slide(prs, spec):
    slide = base_slide(prs)
    add_header(slide, spec["num"], spec["title"], spec.get("subtitle"))

    accent = C.get(spec.get("accent", "neon_blue"), C["neon_blue"])

    # Left panel: KEY POINTS (60% width)
    add_panel(slide, Inches(0.5), Inches(1.6), Inches(8.0), Inches(4.0),
              label="KEY POINTS", accent=accent)
    add_bullets(slide, spec["points"],
                Inches(0.75), Inches(2.0), Inches(7.5), Inches(3.5),
                font_size=14, color=C["text_primary"], bullet_color=accent)

    # Right panel: EVIDENCE
    add_panel(slide, Inches(8.7), Inches(1.6), Inches(4.1), Inches(4.0),
              label="EVIDENCE / REF", accent=C["neon_green"])
    box = slide.shapes.add_textbox(Inches(8.9), Inches(2.0), Inches(3.7), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = spec["evidence"]
    run.font.name = FONT_MONO
    run.font.size = Pt(10)
    run.font.color.rgb = C["text_secondary"]

    # CTA: Question (full width bottom band)
    add_cta_box(slide, spec["question"])
    add_footer(slide)


def render_title_slide(prs):
    """Title slide for the deck."""
    slide = base_slide(prs)
    # Background gradient effect via large panel
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5), C["panel"])
    # Cyan vertical line on left
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(0.12), Inches(3.5), C["neon_blue"])

    # Brand label
    add_text(slide, "MIGHTY SKILL-BRIDGE", Inches(0.8), Inches(2.2),
             Inches(11.5), Inches(0.4),
             font_size=14, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Title
    add_text(slide, "定例打合せアジェンダ", Inches(0.8), Inches(2.7),
             Inches(11.5), Inches(0.9),
             font_size=40, color=C["text_primary"], bold=True, font_name=FONT_HEADER,
             align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(slide, "社内パイロット報告 ＆ 限定デモGo / 一般公開No-Go判定", Inches(0.8), Inches(3.7),
             Inches(11.5), Inches(0.5),
             font_size=16, color=C["neon_green"], font_name=FONT_BODY,
             align=PP_ALIGN.LEFT)

    # Meeting details
    details = "日時: 2026年6月17日 (水) 13:30 - 14:30想定\n対象: 小林雅水 代表取締役社長 ＆ 開発チーム"
    add_text(slide, details, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8),
             font_size=12, color=C["text_secondary"], font_name=FONT_BODY,
             align=PP_ALIGN.LEFT)

    add_footer(slide)


def main():
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Render Title
    render_title_slide(prs)

    # Render Slides
    for spec in SLIDES:
        render_slide(prs, spec)

    prs.save(str(PPTX_FILE))
    print(f"[+] Branded CEO Meeting Agenda PPTX generated successfully: {PPTX_FILE}")


if __name__ == "__main__":
    main()
