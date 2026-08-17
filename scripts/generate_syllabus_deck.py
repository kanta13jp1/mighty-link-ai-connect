#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate 60-Minute Antigravity Live Demo Syllabus PowerPoint Deck (Canva/PPTX).

Event Date: 2026-08-26
Theme: 60-Minute Antigravity Live Demo & 6-Step Workflow Syllabus
Format: 16:9 Wide (1920x1080 / 13.333x7.5 in)

Syllabus Slides Structure:
1. Cover: 研修シラバス - Antigravity 60分ライブデモ（Web生成からデプロイ・ガードまで）
2. Course Overview & Target: 講座概要・対象者・ゴール・前提知識
3. 60-Min Timetable: 全体スケジュール（導入10分 / 6大実演30分 / ハンズオン10分 / Q&A 10分）
4. 6-Step Live Demo Curriculum: 6大実演ステップの詳細シラバス（Step 1〜3）
5. 6-Step Live Demo Curriculum: 6大実演ステップの詳細シラバス（Step 4〜6）
6. Participant Prep & Environment: 参加者事前準備・動作環境・配布教材
7. Safety Guidelines & Rules: ガバナンス・セキュリティ・ファクトチェック
8. Post-Training Roadmap & Support: 研修後のフォローアップ・Slackサポート窓口
"""

from __future__ import annotations

import datetime as dt
import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports" / "training_deck"
PPTX_FILE = EXPORT_DIR / "antigravity_60min_demo_syllabus_2026-08-26.pptx"
SUMMARY_FILE = EXPORT_DIR / "antigravity_60min_demo_syllabus_2026-08-26.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Canva Clean Studio Theme Palette (Mint & Slate)
COLORS = {
    "bg_dark": RGBColor(11, 15, 25),          # Deep Canvas #0B0F19
    "card_dark": RGBColor(21, 30, 46),        # Surface Slate #151E2E
    "card_elevated": RGBColor(30, 41, 59),    # Highlight Container #1E293B
    "border_subtle": RGBColor(42, 54, 79),    # Divider line #2A364F
    
    # Canva Mint & Gradient Accents
    "mint_bright": RGBColor(45, 212, 191),    # Teal/Mint 400 #2DD4BF
    "mint_bg": RGBColor(13, 53, 50),          # Mint Deep Container
    "sky_cyan": RGBColor(56, 189, 248),       # Sky 400 #38BDF8
    "purple_accent": RGBColor(167, 139, 250), # Purple 400 #A78BFA
    "amber_accent": RGBColor(251, 191, 36),   # Amber 400 #FBBF24
    "emerald_accent": RGBColor(52, 211, 153), # Emerald 400 #34D399
    "rose_accent": RGBColor(251, 113, 133),   # Rose 400 #FB7185
    
    # Text
    "text_title": RGBColor(248, 250, 252),    # #F8FAFC
    "text_body": RGBColor(203, 213, 225),     # #CBD5E1
    "text_muted": RGBColor(148, 163, 184),    # #94A3B8
    "text_subtle": RGBColor(100, 116, 139),   # #64748B
}

FONT_DISPLAY = "Plus Jakarta Sans"
FONT_BODY = "Yu Gothic"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name or FONT_BODY
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color or COLORS["text_title"]
    return box


def add_bullets(
    slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 13,
    color: RGBColor | None = None,
    font_name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"✔  {bullet}"
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = font_name or FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["text_body"]
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_chip(
    slide,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    text: RGBColor,
    border: RGBColor | None = None,
):
    add_rect(slide, x, y, w, 0.34, fill=fill, line=border or fill, radius=True)
    add_text(
        slide,
        label,
        x,
        y + 0.04,
        w,
        0.24,
        size=9.5,
        color=text,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        font_name=FONT_DISPLAY,
    )


def add_header(slide, slide_no: int, eyebrow: str, title: str):
    add_rect(slide, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], line=COLORS["bg_dark"])
    
    # Canva Dual Accent Bar
    add_rect(slide, 0.75, 0.35, 0.6, 0.05, fill=COLORS["mint_bright"], line=COLORS["mint_bright"])
    add_rect(slide, 1.4, 0.35, 0.35, 0.05, fill=COLORS["sky_cyan"], line=COLORS["sky_cyan"])
    
    # Eyebrow
    add_text(
        slide,
        f"{slide_no:02d}  /  {eyebrow.upper()}",
        0.75,
        0.48,
        8.0,
        0.26,
        size=9,
        color=COLORS["mint_bright"],
        bold=True,
        font_name=FONT_MONO,
    )
    # Title
    add_text(
        slide,
        title,
        0.73,
        0.76,
        11.8,
        0.65,
        size=25,
        color=COLORS["text_title"],
        bold=True,
        font_name=FONT_DISPLAY,
    )
    
    # Footer
    add_text(
        slide,
        f"2026-08-26 社内AI研修シラバス · Antigravity 60分完全実演  ·  Page {slide_no:02d} / 08",
        0.75,
        7.05,
        7.2,
        0.28,
        size=8.5,
        color=COLORS["text_subtle"],
        font_name=FONT_MONO,
    )
    add_text(
        slide,
        "Mighty Link AI Connect  ·  Training Syllabus",
        8.0,
        7.05,
        4.58,
        0.28,
        size=8.5,
        color=COLORS["text_subtle"],
        align=PP_ALIGN.RIGHT,
        font_name=FONT_DISPLAY,
    )


def build_syllabus_deck(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]

    # --- Slide 1: Syllabus Cover ---
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=COLORS["bg_dark"], line=COLORS["bg_dark"])
    
    add_rect(slide, 0.75, 0.85, 11.83, 5.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    add_chip(slide, "研修シラバス (SYLLABUS) · 2026-08-26 開催", 1.2, 1.25, 4.0, fill=COLORS["mint_bg"], text=COLORS["mint_bright"], border=COLORS["mint_bright"])
    
    add_text(slide, "Google Antigravity 60分ライブデモ実践研修", 1.15, 1.75, 11.0, 0.65, size=29, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "Web生成 ➜ 画面直接修正 ➜ 自律検証 ➜ 公開 ➜ 健全性ガード ➜ ローカル連携", 1.15, 2.45, 11.0, 0.55, size=20, color=COLORS["mint_bright"], bold=True, font_name=FONT_DISPLAY)
    add_text(slide, "〜 生成AIチャットの壁を超え、AIエージェントと1時間でWebアプリを共創・公開する標準カリキュラム 〜", 1.18, 3.15, 11.0, 0.35, size=14, color=COLORS["text_body"])

    meta_boxes = [
        ("受講対象", "全社（非エンジニア・企画・営業・開発者）", COLORS["mint_bright"]),
        ("所要時間", "オンライン 60分（実演＋体験）", COLORS["sky_cyan"]),
        ("到達目標", "自然言語での要件定義・画面修正指示を体得", COLORS["purple_accent"]),
    ]
    for idx, (label, val, accent) in enumerate(meta_boxes):
        x = 1.18 + idx * 3.7
        add_rect(slide, x, 3.75, 3.45, 2.5, fill=COLORS["card_elevated"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.2, 4.0, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, label, x + 0.2, 4.2, 3.05, 0.3, size=14, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, val, x + 0.2, 4.65, 3.05, 1.3, size=12, color=COLORS["text_muted"])

    # --- Slide 2: Course Overview & Target ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 2, "Course Overview", "講座のねらいと受講対象・到達目標")
    add_text(slide, "プログラミング経験不問。AIエージェントに「正しく指示を出し、安全に使いこなす」スキルを習得します。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    cards_overview = [
        ("🎯 研修のねらい", "「ChatGPTで文章を作る」段階から、「AIエージェントにアプリ・サイトを作らせて業務を自動化する」次世代の仕事術を60分で体感する。", COLORS["mint_bright"]),
        ("👥 対象受講者", "・新規Web施策やLPを素早く作りたい企画・営業\n・社内ツールやマニュアルを自動化したい非エンジニア\n・AIエージェントの最新開発フローを学びたいエンジニア", COLORS["sky_cyan"]),
        ("🏆 研修後の到達目標", "① 自然言語でAIに過不足のない要件を伝えられる\n② 画面プレビューを見ながらピン留め修正指示ができる\n③ 自動公開された本番サイトの健全性を検証できる", COLORS["emerald_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(cards_overview):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, 2.2, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, 2.45, 3.3, 0.4, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.25, 3.05, 3.3, 3.4, size=12.5, color=COLORS["text_body"])

    # --- Slide 3: 60-Min Timetable ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 3, "Schedule & Structure", "60分間のタイムテーブル構成")
    add_text(slide, "だれもが飽きずに集中できる「解説 ➜ 生実演 ➜ 体験 ➜ 質疑」の最適配分", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    blocks = [
        ("10分", "01. 導入と背景", "・AIエージェントの現在地 (2026年)\n・ChatGPT SitesとAntigravityの違い\n・本日の実演テーマの発表", COLORS["mint_bright"]),
        ("30分", "02. 6大ステップ完全生実演", "・Webサイト生成からローカル連携まで\n・画面直接ピン留め修正のノーカット実演\n・リアルタイム本番デプロイ", COLORS["sky_cyan"]),
        ("10分", "03. 参加者プチ体験", "・公開URLを参加者自身の端末で閲覧\n・Visual Feedbackのピン留め指示体験\n・スマホ表示でのレスポンシブ確認", COLORS["purple_accent"]),
        ("10分", "04. 質疑応答 & 申請案内", "・現場の疑問に答えるリアルタイムQ&A\n・社内アカウント申請フォーム案内\n・社内プロンプト共有リポジトリ紹介", COLORS["emerald_accent"]),
    ]
    for idx, (time_tag, title, desc, accent) in enumerate(blocks):
        x = 0.75 + idx * 3.01
        add_rect(slide, x, 1.95, 2.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, time_tag, x + 0.2, 2.2, 1.2, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, x + 0.2, 2.75, 2.4, 0.55, size=15, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.2, 3.45, 2.4, 3.0, size=12, color=COLORS["text_body"])

    # --- Slide 4: 6-Step Curriculum Part 1 (Steps 1-3) ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 4, "Curriculum: Steps 1 - 3", "実演カリキュラム詳細（ステップ 1 〜 3）")
    add_text(slide, "要件定義からAIの自動実装、画面上での直感的な修正指示までを解説・実演します。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    steps_p1 = [
        ("STEP 01", "Webサイト生成", "プロンプト投入 ➜ 計画策定 ➜ コード生成", [
            "自然言語で「9大AIエージェント比較ポータルを作って」と指示",
            "AntigravityがImplementation Plan（設計書）を自律立案",
            "HTML/CSS/JSを自動生成し右側Artifactsにプレビュー即時表示",
        ], COLORS["mint_bright"]),
        ("STEP 02", "修正・レビュー体験", "Visual Feedbackによる画面直接ピン留め指示", [
            "プレビュー画面上のボタンや見出しを直接マウスクリック",
            "「📌 ピン」を立てて「ここをミント色にして」と短い指示を送信",
            "長い文章の説明が不要になり、指示ミス・手戻りがゼロに",
        ], COLORS["sky_cyan"]),
        ("STEP 03", "ブラウザ自動検証", "Browser Agentによる自律E2Eテスト", [
            "Headless Chrome / Playwright をAIがバックグラウンド起動",
            "スマホ画面（390px）とPC画面（1440px）の崩れを自動検査",
            "ボタンクリック・モーダル表示の挙動をAIが自己修復",
        ], COLORS["purple_accent"]),
    ]
    for idx, (step_num, title, sub, bullets, accent) in enumerate(steps_p1):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, step_num, x + 0.25, 2.2, 1.4, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, x + 0.25, 2.68, 3.3, 0.35, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, sub, x + 0.25, 3.1, 3.3, 0.45, size=11, color=accent)
        add_bullets(slide, bullets, x + 0.25, 3.65, 3.3, 2.8, size=11.5)

    # --- Slide 5: 6-Step Curriculum Part 2 (Steps 4-6) ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 5, "Curriculum: Steps 4 - 6", "実演カリキュラム詳細（ステップ 4 〜 6）")
    add_text(slide, "世界中からアクセス可能な公開、自動ガードレール検証、ローカル開発環境との同期を実演します。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    steps_p2 = [
        ("STEP 04", "ホスティング・公開", "GitHub Pages / Firebaseへの自動デプロイ", [
            "「本番公開へデプロイして」と指示するだけでCI/CDが起動",
            "GitHub Actions経由でHTTPS対応の本番URLが即座に発行",
            "ChatGPT Sites同様のワンクリック公開を自社インフラで実現",
        ], COLORS["emerald_accent"]),
        ("STEP 05", "公開URLの健全性ガード", "自動テストスクリプトによるデグレ防止", [
            "デプロイ完了後、AIが自動で健全性ガード（E2E疎通）を実行",
            "HTTP 200ステータス、全リンク到達性、DOM要素存在を検査",
            "「壊れたサイトが勝手に公開される事故」を完全に防止",
        ], COLORS["amber_accent"]),
        ("STEP 06", "コードのローカル連携", "Gitリポジトリ・VS Codeとの完全同期", [
            "Web完結ツールと異なり、生成コードが手元PCにそのまま残る",
            "既存の社内GitリポジトリやCI/CDパイプラインに即時統合可能",
            "ブラックボックス化せず、開発チームへスムーズに引き継げる",
        ], COLORS["rose_accent"]),
    ]
    for idx, (step_num, title, sub, bullets, accent) in enumerate(steps_p2):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, step_num, x + 0.25, 2.2, 1.4, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, x + 0.25, 2.68, 3.3, 0.35, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, sub, x + 0.25, 3.1, 3.3, 0.45, size=11, color=accent)
        add_bullets(slide, bullets, x + 0.25, 3.65, 3.3, 2.8, size=11.5)

    # --- Slide 6: Participant Prep & Environment ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 6, "Requirements & Prep", "受講に必要な事前準備・推奨環境")
    add_text(slide, "参加者の皆様は特別なツールインストールの必要なく、ブラウザのみでご参加いただけます。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    prep_cards = [
        ("💻 受講環境（参加者）", [
            "Google Chrome / Microsoft Edge（最新版ブラウザ）",
            "オンライン会議ツール（Zoom / Google Meet / Teams）",
            "スマートフォン（レスポンシブ表示確認用・任意）",
            "※ 特別の開発環境構築や拡張機能インストールは不要です",
        ], COLORS["mint_bright"]),
        ("📦 配布教材・資料", [
            "本研修スライド（PDF / PPTX / Canva共有リンク）",
            "実演で使用するプロンプト定石テンプレート集",
            "実演公開デモURL（https://kanta13jp1.github.io/...）",
            "Antigravity社内利用スタートアップマニュアル",
        ], COLORS["sky_cyan"]),
        ("⚙️ 講師側実演環境", [
            "Google Antigravity IDE & Gemini 3.7 Flash Engine",
            "GitHub Repository / GitHub Actions CI/CD Pipeline",
            "Headless Chrome E2E 自動検証スクリプト",
            "Canva & Figma REST API 連携環境",
        ], COLORS["purple_accent"]),
    ]
    for idx, (title, bullets, accent) in enumerate(prep_cards):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 4.8, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, 2.2, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, 2.45, 3.3, 0.4, size=16, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_bullets(slide, bullets, x + 0.25, 3.1, 3.3, 3.4, size=12)

    # --- Slide 7: Security & Governance ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 7, "Security & Rules", "社内利用における安全ガイドラインと禁止事項")
    add_text(slide, "安心・安全にAIエージェントを全社展開するために遵守すべき4つのルール", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    rules = [
        ("RULE 01: 機密情報・個人情報の投入禁止", "顧客の実名、パスワード、APIキー、未公開の社内規程・契約書は絶対にプロンプトに入力しない。ダミーデータを使用する。", COLORS["rose_accent"]),
        ("RULE 02: 生成物のファクトチェック・テスト必須", "AIの出力には誤情報（ハルシネーション）が含まれうるため、生成されたコードや文章は必ず人間が動作確認・精査する。", COLORS["amber_accent"]),
        ("RULE 03: 外部公開時の著作権・ライセンス確認", "生成した画像・アイコン・文章を社外公開する際は、商用利用規約およびOSSライセンス条項に抵触しないことを確認する。", COLORS["sky_cyan"]),
        ("RULE 04: クォータとコストの適正管理", "モデルの利用上限（クォータ）を意識し、大量バッチ処理や連続実行を行う際は事前に管理部門（情シス）へ相談する。", COLORS["emerald_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(rules):
        row_y = 1.95 + idx * 1.18
        add_rect(slide, 0.75, row_y, 11.83, 1.05, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_chip(slide, f"SECURITY", 0.95, row_y + 0.15, 1.3, fill=COLORS["card_elevated"], text=accent, border=COLORS["border_subtle"])
        add_text(slide, title, 2.45, row_y + 0.15, 9.8, 0.3, size=13.5, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, 2.45, row_y + 0.48, 9.8, 0.45, size=11, color=COLORS["text_body"])

    # --- Slide 8: Next Steps & Syllabus Summary ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, 8, "Roadmap & Support", "研修受講後の実践ロードマップとサポート")
    add_text(slide, "受講して終わりにしない！各部署での実務適用とスキル定着をサポートします。", 0.78, 1.5, 11.7, 0.35, size=14, color=COLORS["text_body"])

    steps = [
        ("STEP 1: アカウント発行", "社内専用申請フォームよりAntigravityアカウントを発行（申請から即日利用可能）", COLORS["mint_bright"]),
        ("STEP 2: 定石プロンプトで実践", "社内共有リポジトリの「業務別プロンプト・Skills定義」を使って自分の業務を1つ自動化", COLORS["sky_cyan"]),
        ("STEP 3: 成果共有と勉強会", "月次開催の「AI活用ライトニングトーク会」で自作したWebサイトや自動化成果を発表", COLORS["purple_accent"]),
    ]
    for idx, (title, desc, accent) in enumerate(steps):
        x = 0.75 + idx * 4.02
        add_rect(slide, x, 1.95, 3.8, 3.2, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
        add_rect(slide, x + 0.25, 2.2, 0.4, 0.05, fill=accent, line=accent)
        add_text(slide, title, x + 0.25, 2.45, 3.3, 0.35, size=15, color=COLORS["text_title"], bold=True, font_name=FONT_DISPLAY)
        add_text(slide, desc, x + 0.25, 3.0, 3.3, 1.8, size=12, color=COLORS["text_body"])

    # Bottom Contact Box
    add_rect(slide, 0.75, 5.35, 11.83, 1.45, fill=COLORS["card_dark"], line=COLORS["border_subtle"], radius=True)
    add_chip(slide, "💬 研修事務局 & 社内相談窓口", 0.95, 5.52, 2.8, fill=COLORS["mint_bg"], text=COLORS["mint_bright"], border=COLORS["mint_bright"])
    add_text(slide, "社内Slack: #ai-agent-hub  |  お問い合わせ: ai-support@ml-mightylink.com  |  次回応用編ワークショップ: 2026年9月中旬予定", 0.95, 6.0, 11.4, 0.5, size=11.5, color=COLORS["text_body"])


def verify_pptx(path: Path, expected_slides: int) -> None:
    if not path.exists() or path.stat().st_size < 10_000:
        raise RuntimeError(f"PPTX was not created or is too small: {path}")
    with zipfile.ZipFile(path) as package:
        slide_parts = [
            name for name in package.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
    if len(slide_parts) != expected_slides:
        raise RuntimeError(f"Expected {expected_slides} slides, found {len(slide_parts)}")


def write_summary() -> None:
    summary_content = f"""# 2026-08-26 Antigravity 60分ライブデモ研修 シラバス

生成日時: {jst_now().isoformat(timespec="seconds")}

## 成果物ファイル
- **PowerPointファイル**: `exports/training_deck/{PPTX_FILE.name}`
- **シラバスドキュメント**: `exports/training_deck/{SUMMARY_FILE.name}`

## シラバス構成 (全8枚 / 16:9 ワイド)
1. **シラバス表紙**: Google Antigravity 60分ライブデモ実践研修（Web生成からデプロイ・ガードまで）
2. **講座概要 & 対象者**: 講座のねらい・受講対象者（非エンジニア含む）・到達目標
3. **60分タイムテーブル**: 導入10分 / 6大ステップ生実演30分 / プチ体験10分 / 質疑応答10分
4. **カリキュラム詳細 (前編)**:
   - Step 1: Webサイト生成（要件投入 ➜ 計画 ➜ 自動コーディング）
   - Step 2: 修正・レビュー体験（Visual Feedback 画面直接ピン留め指示）
   - Step 3: ブラウザ自動検証（Browser Agent によるE2E自動テスト & レスポンシブ検査）
5. **カリキュラム詳細 (後編)**:
   - Step 4: ホスティング・公開（GitHub Pages / Firebaseへの自動デプロイ）
   - Step 5: 公開URLの健全性ガード（HTTP 200・リンク切れ0件の自動監査）
   - Step 6: コードのローカル連携（Gitリポジトリ・VS Codeとの完全同期）
6. **受講環境 & 事前準備**: 参加者環境（ブラウザのみ）・配布教材・講師実演環境
7. **セキュリティ & 社内ルール**: 機密情報投入禁止・ファクトチェック・著作権・クォータ管理
8. **研修後の実践ロードマップ**: アカウント発行・社内プロンプト集・Slack相談窓口
"""
    SUMMARY_FILE.write_text(summary_content, encoding="utf-8")


def main() -> None:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    build_syllabus_deck(prs)
    prs.save(PPTX_FILE)
    
    verify_pptx(PPTX_FILE, expected_slides=8)
    write_summary()
    
    print("[+] Antigravity 60分ライブデモ研修 シラバスPPTX (全8枚) を生成しました。")
    print(f"[*] PPTX: {PPTX_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary: {SUMMARY_FILE.relative_to(PROJECT_ROOT)}")


if __name__ == "__main__":
    main()
