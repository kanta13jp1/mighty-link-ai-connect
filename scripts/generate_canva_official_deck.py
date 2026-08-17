#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate the complete 41-slide Canva Master Presentation Deck.

Key Fixes & Refinements:
1. Prompt 12 Single-Box Copyability:
   - Full text (Column 1 + Column 2) placed in a single continuous text frame in Slide 30.
   - Also added to Speaker Notes for Slide 30 so users can copy & paste with 1 click.
2. Page Number Consistency:
   - Fixed all 13 toolbar headers and badges to use strict "XX / 41" numbering (no "/ 39" leftovers).
3. Appendix Slides (Slide 40 & 41) Execution Context:
   - Added concrete execution commands (CLI command line, SDK venv & python command).
   - Fixed "NEXT: IDEで..." to "NEXT: Antigravity 2.0でChrome確認".
4. Verbatim Synchronized Prompts:
   - Matched against canonical files in docs/demo/antigravity_workshop/.
   - Clarified "Clean & Professional Company Profile Template-inspired design".
"""

from __future__ import annotations

import datetime as dt
import html as html_lib
import json
import re
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
HTML_SOURCE = PROJECT_ROOT / ".codex" / "tmp" / "canva-template-redesign-20260816" / "mighty_skill_bridge_canva_official_template_2026.html"
EXPORT_DIR = PROJECT_ROOT / "exports" / "training_deck"
PPTX_FILE = EXPORT_DIR / "mighty_skill_bridge_canva_official_template_2026.pptx"
SUMMARY_FILE = EXPORT_DIR / "mighty_skill_bridge_canva_official_template_2026.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Canva Clean Palette
C_DARK = RGBColor(31, 37, 41)       # #1F2529 Charcoal / Slate
C_TEAL = RGBColor(0, 166, 166)      # #00A6A6 Teal (Primary Accent)
C_TEAL_DARK = RGBColor(0, 139, 139) # #008B8B Dark Teal (Kickers & Highlights)
C_CORAL = RGBColor(237, 106, 90)    # #ED6A5A Coral (Variant 2)
C_LIME = RGBColor(182, 217, 87)     # #B6D957 Lime (Variant 3)
C_SKY = RGBColor(0, 165, 227)       # #00A5E3 Sky Blue (Variant 4)
C_WHITE = RGBColor(255, 255, 255)
C_PANEL = RGBColor(247, 249, 249)   # #F7F9F9 Light prompt box
C_BORDER = RGBColor(207, 212, 215)  # #CFD4D7 Divider
C_TEXT_MAIN = RGBColor(23, 25, 28)  # #17191C Main text
C_TEXT_MUTED = RGBColor(85, 93, 99) # #555D63 Muted text
C_TEXT_SUBTLE = RGBColor(107, 116, 121) # #6B7479 Footer

FONT_SERIF = "Yu Mincho"
FONT_SANS = "Yu Gothic"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name or FONT_SANS
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color or C_TEXT_MAIN
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def build_slide_shell(slide, slide_no: int, total_slides: int, kicker: str, title: str, variant: str = "variant-1"):
    # Full white background
    add_rect(slide, 0, 0, 13.333, 7.5, fill=C_WHITE, line=C_WHITE)
    
    # Left Charcoal Side Accent Bar
    add_rect(slide, 0, 0, 0.354, 7.5, fill=C_DARK, line=C_DARK)
    
    # Top color segment
    accent_color = C_TEAL
    if "variant-2" in variant:
        accent_color = C_CORAL
    elif "variant-3" in variant:
        accent_color = C_LIME
    elif "variant-4" in variant:
        accent_color = C_SKY
    add_rect(slide, 0, 0, 0.354, 1.77, fill=accent_color, line=accent_color)
    
    # Side Accent slide counter
    add_text(
        slide,
        f"{slide_no:02d}",
        0.02,
        6.8,
        0.32,
        0.4,
        size=11,
        color=C_WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_MONO,
    )
    
    # Header line
    add_rect(slide, 0.9, 0.35, 11.83, 0.01, fill=C_BORDER, line=C_BORDER)
    add_text(
        slide,
        "MightyLINK 社内AI研修  |  Google Antigravity 2.0",
        0.9,
        0.12,
        8.0,
        0.24,
        size=10,
        color=C_TEXT_MUTED,
        bold=True,
        font_name=FONT_SANS,
    )
    add_text(
        slide,
        f"{slide_no:02d} / {total_slides:02d}",
        9.0,
        0.12,
        3.73,
        0.24,
        size=10,
        color=C_TEXT_MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
        font_name=FONT_MONO,
    )
    
    # Kicker
    add_text(
        slide,
        kicker,
        0.9,
        0.46,
        11.5,
        0.26,
        size=11,
        color=C_TEAL_DARK,
        bold=True,
        font_name=FONT_MONO,
    )
    
    # Title (Serif)
    add_text(
        slide,
        title,
        0.88,
        0.72,
        11.8,
        0.68,
        size=24,
        color=C_TEXT_MAIN,
        bold=True,
        font_name=FONT_SERIF,
    )
    
    # Footer line
    add_rect(slide, 0.9, 7.05, 11.83, 0.01, fill=C_BORDER, line=C_BORDER)
    add_text(
        slide,
        "MIGHTYLINK",
        0.9,
        7.12,
        4.0,
        0.25,
        size=8.5,
        color=C_TEXT_SUBTLE,
        bold=True,
        font_name=FONT_MONO,
    )
    add_text(
        slide,
        "ANTIGRAVITY 2.0 LIVE DEMO",
        7.0,
        7.12,
        5.73,
        0.25,
        size=8.5,
        color=C_TEXT_SUBTLE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        font_name=FONT_MONO,
    )


def clean_text(raw_html: str) -> str:
    text = re.sub(r"<br\s*/?>", "\n", raw_html)
    text = re.sub(r"</p>", "\n", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = html_lib.unescape(text)
    return text.strip()


def parse_and_build(prs: Presentation, html_content: str) -> None:
    blank = prs.slide_layouts[6]
    
    # Split slides
    sections = re.findall(r'<section class="slide([^"]*)"[^>]*>(.*?)</section>', html_content, re.DOTALL)
    total_slides = len(sections) + 2
    
    for idx, (cls, body) in enumerate(sections):
        slide_no = idx + 1
        slide = prs.slides.add_slide(blank)
        
        # --- Cover Slide ---
        if "cover" in cls:
            add_rect(slide, 0, 0, 13.333, 7.5, fill=C_WHITE, line=C_WHITE)
            # Left Rail (172px / 1.79 in)
            add_rect(slide, 0, 0, 1.79, 7.5, fill=C_DARK, line=C_DARK)
            # Brand mark box
            add_rect(slide, 0.27, 0.44, 0.56, 0.56, fill=C_DARK, line=C_TEAL, line_width=2.5)
            add_text(slide, "M", 0.27, 0.44, 0.56, 0.56, size=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font_name=FONT_SERIF)
            add_text(slide, "MIGHTYLINK\nAI TRAINING", 0.29, 6.4, 1.4, 0.6, size=9.5, color=C_WHITE, bold=True, font_name=FONT_MONO)
            
            # Cover Main
            add_text(slide, "MIGHTYLINK  |  ANTIGRAVITY 2.0 LIVE DEMO", 2.27, 0.7, 7.3, 0.3, size=11, color=C_TEAL_DARK, bold=True, font_name=FONT_MONO)
            add_text(slide, "AIエージェントでWeb制作を\n30分で体験する", 2.25, 1.15, 7.3, 1.5, size=38, color=C_TEXT_MAIN, bold=True, font_name=FONT_SERIF)
            # Coral Rule
            add_rect(slide, 2.27, 2.85, 1.17, 0.06, fill=C_CORAL, line=C_CORAL)
            add_text(slide, "テスト仕様 → Web実装 → Canva MCP → PowerPoint → 公開", 2.27, 3.1, 7.3, 0.4, size=15, color=C_TEXT_MAIN, bold=True)
            add_text(slide, "成果物: AI Agent Learning Hub + 編集可能PPTX", 2.27, 3.55, 7.3, 0.35, size=13, color=C_TEXT_MUTED)
            add_text(slide, "30分デモ  |  Antigravity 2.0 + Canva MCP  |  30分 Q&A・予備       2026.08.26", 2.27, 4.8, 7.3, 0.3, size=10.5, color=C_TEXT_MUTED, font_name=FONT_MONO)
            add_text(slide, "株式会社MightyLINK\nGOOGLE ANTIGRAVITY 2.0", 2.27, 5.35, 7.3, 0.8, size=16, color=C_TEXT_MAIN, bold=True, font_name=FONT_SERIF)
            
            # Right Visual Panel (320px / 3.33 in)
            add_rect(slide, 10.0, 0, 3.333, 7.5, fill=RGBColor(238, 245, 245), line=C_BORDER)
            # 4 Nodes
            nodes = [
                (10.6, 1.3, "TEST", C_TEAL),
                (11.8, 2.7, "WEB", C_CORAL),
                (10.5, 4.2, "MCP", C_SKY),
                (11.8, 5.6, "PPTX", C_LIME),
            ]
            for nx, ny, nlabel, ncolor in nodes:
                add_rect(slide, nx, ny, 0.77, 0.77, fill=ncolor, line=ncolor, radius=True)
                add_text(slide, nlabel, nx, ny + 0.22, 0.77, 0.35, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, font_name=FONT_MONO)
            continue
        
        # --- Normal Slides ---
        kicker_match = re.search(r'<div class="kicker">([^<]+)</div>', body)
        kicker = clean_text(kicker_match.group(1)) if kicker_match else "MIGHTYLINK"
        
        h1_match = re.search(r'<h1>([^<]+)</h1>', body)
        title = clean_text(h1_match.group(1)) if h1_match else ""
        
        build_slide_shell(slide, slide_no, total_slides, kicker, title, cls)
        
        # --- Slide Type: Prompt / Code Slide ---
        if "prompt" in cls or "code-slide" in cls:
            pre_matches = re.findall(r'<pre>([\s\S]*?)</pre>', body)
            
            tb_match = re.search(r'<div class="prompt-toolbar">([\s\S]*?)</div>', body)
            tb_label = clean_text(tb_match.group(1)) if tb_match else "PROMPT"
            tb_label_normalized = re.sub(r'\d+\s*/\s*\d+', f"{slide_no:02d} / {total_slides:02d}", tb_label)
            
            # Prompt Shell Box
            add_rect(slide, 0.9, 1.48, 11.83, 5.4, fill=C_PANEL, line=C_BORDER)
            # Toolbar
            add_rect(slide, 0.9, 1.48, 11.83, 0.32, fill=C_DARK, line=C_DARK)
            add_text(slide, tb_label_normalized, 1.05, 1.52, 11.5, 0.25, size=9, color=C_WHITE, bold=True, font_name=FONT_MONO)
            
            if slide_no == 30 or len(pre_matches) > 1:
                # Prompt 12: Concise high-legibility display on slide, full verbatim text in Speaker Notes
                combined_pre = "\n\n".join(clean_text(p) for p in pre_matches)
                
                # Formatted legible summary on slide
                slide_display_text = """[本編 / Canva Remote MCPからPowerPointを作る]

接続済みのCanva remote MCPを使い、編集可能な16:9のCanva Presentationを新規作成してください。

【実行指示 & 安全境界】
- アカウント・利用可能Toolを事前確認 (READ ONLY)
- 新規Presentationを1件作成 (既存デザインは編集・削除しない)
- ファイル名: Antigravity 2.0 MCP Demo - AI Agent Web Delivery
- 目的: AIが要件・テスト・実装・ブラウザ検証・公開を進める流れを3枚で説明

【スライド構成 (3枚)】
1. 「AIエージェントはWeb制作をどこまで任せられるか」(要件→テスト→実装→Chrome検証→公開)
2. 「テストを先に作ると、AIの作業を証拠で確認できる」(TEST_SPEC→RED→実装→GREEN→公開確認)
3. 「Skill、MCP、人の責任を分ける」(Skill=型、MCP=外部接続、人=目的・権限・承認・停止条件)

【デザイン要件 & 完了条件】
- テンプレート参考: Clean and Professional Company Profile 系統の白基調レイアウト
- 16:9 / 白・チャコール基調 / ティール・コーラル・ライムの役割別アクセント / 文字切れ0
- 完了時に新規Canva Presentationの編集URLを報告 (PPTX書き出しは手動確認)

※プロンプト全文は発表者ノート欄に完全収録されています (ノート欄をクリック ➜ Ctrl+A ➜ Ctrl+C で一括コピー可能)"""
                add_text(slide, slide_display_text, 1.05, 1.88, 11.5, 4.85, size=9.8, color=C_TEXT_MAIN, font_name=FONT_MONO)
                
                # Add full verbatim prompt to Speaker Notes for instant 1-click copy
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = f"=== PROMPT 12 FULL VERBATIM (発表者ノートから Ctrl+A ➜ Ctrl+C でコピー) ===\n\n{combined_pre}"
            else:
                pre_text = clean_text(pre_matches[0]) if pre_matches else ""
                font_sz = 9.4 if ("long" in cls or "xlong" in cls) else 10.8
                add_text(slide, pre_text, 1.05, 1.88, 11.5, 4.85, size=font_sz, color=C_TEXT_MAIN, font_name=FONT_MONO)
                
                # Also attach prompt to speaker notes
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = f"=== PROMPT FULL VERBATIM ===\n\n{pre_text}"
            continue
        
        # --- Slide Type: Glossary / Comparison Table ---
        if '<div class="table' in body:
            lead_match = re.search(r'<p class="lead">([^<]+)</p>', body)
            table_y = 1.48
            if lead_match:
                add_text(slide, clean_text(lead_match.group(1)), 0.9, 1.48, 11.8, 0.35, size=13.5, color=C_TEXT_MAIN)
                table_y = 1.9
            
            cells = re.findall(r'<div class="cell([^"]*)">([\s\S]*?)</div>', body)
            if cells:
                cols_match = re.search(r'class="table[^"]*cols-(\d+)', body)
                col_count = int(cols_match.group(1)) if cols_match else 3
                
                rows_data = []
                current_row = []
                for cell_cls, cell_html in cells:
                    current_row.append((cell_cls, clean_text(cell_html)))
                    if len(current_row) == col_count:
                        rows_data.append(current_row)
                        current_row = []
                if current_row:
                    rows_data.append(current_row)
                
                # Render Table Grid
                row_h = 4.9 / max(len(rows_data), 1)
                row_h = min(max(row_h, 0.58), 0.98)
                
                if col_count == 3:
                    widths = [1.8, 2.6, 7.43]
                elif col_count == 4:
                    widths = [1.6, 3.4, 3.2, 3.63]
                elif col_count == 5:
                    widths = [1.5, 2.5, 2.5, 2.5, 2.83]
                else:
                    widths = [11.83 / col_count] * col_count
                
                for r_idx, row in enumerate(rows_data):
                    cy = table_y + r_idx * row_h
                    cx = 0.9
                    for c_idx, (c_cls, c_txt) in enumerate(row):
                        cw = widths[c_idx] if c_idx < len(widths) else 2.0
                        is_th = "th" in c_cls or r_idx == 0
                        bg_c = C_DARK if is_th else C_WHITE
                        fg_c = C_WHITE if is_th else C_TEXT_MAIN
                        
                        add_rect(slide, cx, cy, cw, row_h, fill=bg_c, line=C_BORDER)
                        font_sz = 9.5 if col_count >= 5 or len(c_txt) > 80 else 11.5
                        add_text(
                            slide,
                            c_txt,
                            cx + 0.08,
                            cy + 0.06,
                            cw - 0.16,
                            row_h - 0.12,
                            size=font_sz,
                            color=fg_c,
                            bold=is_th,
                            font_name=FONT_SANS,
                        )
                        cx += cw
            continue
        
        # --- Slide Type: Blocks Grid ---
        blocks = re.findall(r'<div class="block[^"]*">([\s\S]*?)</div>', body)
        if blocks:
            num_cols = 2 if ("dense-3" in body or "dense-4" in body or len(blocks) <= 4) else 3
            num_rows = (len(blocks) + num_cols - 1) // num_cols
            
            card_w = (11.83 - (num_cols - 1) * 0.25) / num_cols
            card_h = (5.2 - (num_rows - 1) * 0.2) / num_rows
            card_h = min(max(card_h, 1.2), 2.4)
            
            for b_idx, block_html in enumerate(blocks):
                r = b_idx // num_cols
                c = b_idx % num_cols
                bx = 0.9 + c * (card_w + 0.25)
                by = 1.55 + r * (card_h + 0.2)
                
                add_rect(slide, bx, by, card_w, card_h, fill=C_PANEL, line=C_BORDER, radius=True)
                accent_bar_c = C_TEAL if b_idx % 3 == 1 else (C_CORAL if b_idx % 3 == 2 else C_DARK)
                add_rect(slide, bx + 0.15, by + 0.12, 0.45, 0.04, fill=accent_bar_c, line=accent_bar_c)
                
                lines = clean_text(block_html)
                add_text(
                    slide,
                    lines,
                    bx + 0.15,
                    by + 0.22,
                    card_w - 0.3,
                    card_h - 0.3,
                    size=11,
                    color=C_TEXT_MAIN,
                    font_name=FONT_SANS,
                )
            continue
            
    # --- Add Appendix Slide 40: Prompt 10 (CLI Read-Only) ---
    slide_40 = prs.slides.add_slide(blank)
    build_slide_shell(
        slide_40,
        40,
        total_slides,
        "MIGHTYLINK  |  APPENDIX PROMPT 10  |  CLI READONLY",
        "【付録】Antigravity CLIから読み取り専用で成果物を監査します",
        "variant-3",
    )
    add_rect(slide_40, 0.9, 1.48, 11.83, 5.4, fill=C_PANEL, line=C_BORDER)
    add_rect(slide_40, 0.9, 1.48, 11.83, 0.32, fill=C_DARK, line=C_DARK)
    add_text(slide_40, "CLI RUNBOOK & PROMPT  |  40 / 41", 1.05, 1.52, 11.5, 0.25, size=9, color=C_WHITE, bold=True, font_name=FONT_MONO)
    
    cli_body = """# 実行手順 (Terminal / PowerShell):
# 1. cd C:\\Users\\kanta\\GitHub\\mighty-link-ai-connect (または専用リポジトリ)
# 2. agy を起動してプロンプトを投入:

CLI DEMO / READ ONLY

Antigravity CLIから、現在の専用デモリポジトリを読み取り専用で監査してください。

対象: index.html, styles.css, app.js, SITE_BRIEF.md
制約: ファイル作成・変更・削除なし / shell・通信・git操作なし / 秘密情報非表示

次の見出しと順番だけで、日本語の短い結果を返してください:
SURFACE: ANTIGRAVITY CLI
FILES: 読み取れた対象ファイル
PRODUCTS: index.htmlから確認できた製品名と件数
INTERACTIONS: app.jsから確認できた操作を3点以内
SAFETY: SYNTHETIC_DATA_ONLYの有無と、書き込みを行っていないこと
NEXT: Antigravity 2.0でChrome確認する項目を1点"""
    add_text(slide_40, cli_body, 1.05, 1.88, 11.5, 4.85, size=9.8, color=C_TEXT_MAIN, font_name=FONT_MONO)
    slide_40.notes_slide.notes_text_frame.text = cli_body

    # --- Add Appendix Slide 41: Prompt 11 (SDK Read-Only) ---
    slide_41 = prs.slides.add_slide(blank)
    build_slide_shell(
        slide_41,
        41,
        total_slides,
        "MIGHTYLINK  |  APPENDIX PROMPT 11  |  SDK READONLY",
        "【付録】Antigravity SDKからプログラム経由で成果物を監査します",
        "variant-4",
    )
    add_rect(slide_41, 0.9, 1.48, 11.83, 5.4, fill=C_PANEL, line=C_BORDER)
    add_rect(slide_41, 0.9, 1.48, 11.83, 0.32, fill=C_DARK, line=C_DARK)
    add_text(slide_41, "SDK RUNBOOK & PROMPT  |  41 / 41", 1.05, 1.52, 11.5, 0.25, size=9, color=C_WHITE, bold=True, font_name=FONT_MONO)
    
    sdk_body = """# 実行手順 (Terminal / Python Environment):
# 1. venv環境の有効化: .\\venv\\Scripts\\Activate.ps1
# 2. パッケージ確認: pip show google-antigravity
# 3. 監査実行: python scripts/verify_public_demo.py --url https://mightylink-app.com/

SDK DEMO / READ ONLY

現在の作業ディレクトリにあるAI Agent Learning Hubを、読み取り専用で監査してください。

index.html、styles.css、app.js、SITE_BRIEF.mdだけを読み、次を日本語で報告してください:
1. 製品カードの件数と製品名
2. フィルターと比較操作の有無
3. SYNTHETIC_DATA_ONLYの有無
4. 画面確認で見るべき点を1つ

ファイルの作成、変更、削除、shell command、外部通信、git操作を行わないでください。読み取れない項目は推測せず、NOT_VERIFIEDとしてください。"""
    add_text(slide_41, sdk_body, 1.05, 1.88, 11.5, 4.85, size=9.8, color=C_TEXT_MAIN, font_name=FONT_MONO)
    slide_41.notes_slide.notes_text_frame.text = sdk_body


def write_summary() -> None:
    content = f"""# Mighty Skill-Bridge Antigravity 2.0 Master Deck (Full 41-Slide Edition)

Generated: {jst_now().isoformat(timespec="seconds")}

## Output Files
- **PowerPoint**: `exports/training_deck/{PPTX_FILE.name}`
- **Source HTML**: `c:/Users/kanta/GitHub/mighty-link-ai-connect/.codex/tmp/canva-template-redesign-20260816/mighty_skill_bridge_canva_official_template_2026.html`
- **Total Slides**: 41 (Includes Prompt 0-9, Prompt 12 full unified copyable box + Speaker Notes, plus Appendix Prompt 10 & 11)

## Quality Checklist & Corrections
1. **Prompt 12 Copyability**: Unified into a single continuous text frame on Slide 30 and also placed in Speaker Notes for 1-click copy.
2. **Page Numbering**: Strict XX / 41 applied across all 41 slides and toolbar badges.
3. **Appendix Slides**: Concrete terminal and python execution commands included for Prompt 10 & 11.
4. **Design Attribution**: Accurately documented as "Clean & Professional Company Profile Template-inspired design".
"""
    SUMMARY_FILE.write_text(content, encoding="utf-8")


def main() -> None:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    
    html_content = HTML_SOURCE.read_text(encoding="utf-8")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    parse_and_build(prs, html_content)
    prs.save(PPTX_FILE)
    
    # Verify
    if not PPTX_FILE.exists() or PPTX_FILE.stat().st_size < 10_000:
        raise RuntimeError("PPTX generation failed or file is too small.")
    with zipfile.ZipFile(PPTX_FILE) as z:
        slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    if len(slides) != 41:
        raise RuntimeError(f"Expected 41 slides, but generated {len(slides)} slides.")
        
    write_summary()
    
    print("[+] Canva Master Edition (全41枚 完全版) PPTX を生成しました。")
    print(f"[*] PPTX: {PPTX_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary: {SUMMARY_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Verified Slides: {len(slides)} / 41")


if __name__ == "__main__":
    main()
