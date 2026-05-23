#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate a Canva/Figma-style branded CEO PPTX using python-pptx.

Mighty Skill-Bridge brand palette (Seedance cinematic theme):
- bg: cyber black #0d0e15
- panel: deep navy #161824
- accent1: neon blue #00f0ff
- accent2: neon green #39ff14
- accent3: neon red (risk) #ff3366
- text primary: cool white #f1f5ff
- text secondary: gray white #c5cae0

Source outline: exports/knowledge_flow/notebooklm_ceo_slide_outline.md
Output: exports/knowledge_flow/mighty_skill_bridge_ceo_presentation_2026-06-02_branded.pptx

This is a Claude Code lane variant (Canva MCP オフライン代替). Real Canva/Figma
template usage requires HANDOFF-14b / HANDOFF-17 (Canva MCP) once OAuth is done.
"""

from __future__ import annotations

import datetime as dt
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports" / "knowledge_flow"
PPTX_FILE = EXPORT_DIR / "mighty_skill_bridge_ceo_presentation_2026-06-02_branded.pptx"
SUMMARY_FILE = EXPORT_DIR / "mighty_skill_bridge_ceo_presentation_2026-06-02_branded.md"
MANIFEST_FILE = EXPORT_DIR / "mighty_skill_bridge_ceo_presentation_2026-06-02_branded.json"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Mighty Skill-Bridge brand palette (Seedance cinematic)
C = {
    "bg": RGBColor(0x0D, 0x0E, 0x15),
    "panel": RGBColor(0x16, 0x18, 0x24),
    "panel_alt": RGBColor(0x1F, 0x22, 0x33),
    "neon_blue": RGBColor(0x00, 0xF0, 0xFF),
    "neon_green": RGBColor(0x39, 0xFF, 0x14),
    "neon_red": RGBColor(0xFF, 0x33, 0x66),
    "neon_yellow": RGBColor(0xFF, 0xD7, 0x00),
    "text_primary": RGBColor(0xF1, 0xF5, 0xFF),
    "text_secondary": RGBColor(0xC5, 0xCA, 0xE0),
    "text_muted": RGBColor(0x7A, 0x83, 0x99),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_HEADER = "Yu Gothic UI"
FONT_BODY = "Yu Gothic UI"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h, *,
             font_size=14, color=None, bold=False,
             font_name=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    color = color or C["text_primary"]
    font_name = font_name or FONT_BODY
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def base_slide(prs):
    """Create a blank slide with cyber-black background and corner accent."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Cyber black background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C["bg"])
    # Top-left neon corner accent (5px line, 40% width)
    add_rect(slide, Inches(0), Inches(0), Inches(5.2), Inches(0.04), C["neon_blue"])
    # Bottom-right neon green dot accent
    add_rect(slide, Inches(13.0), Inches(7.2), Inches(0.18), Inches(0.18), C["neon_green"])
    return slide


def add_header(slide, slide_num, title, subtitle=None):
    """Top header with slide number, title, and accent line."""
    # Slide number badge
    add_text(slide, f"0{slide_num}", Inches(0.5), Inches(0.35), Inches(0.8), Inches(0.5),
             font_size=14, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Title
    add_text(slide, title, Inches(1.3), Inches(0.3), Inches(11.5), Inches(0.6),
             font_size=24, color=C["text_primary"], bold=True, font_name=FONT_HEADER,
             align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(1.3), Inches(0.85), Inches(11.5), Inches(0.4),
                 font_size=12, color=C["neon_green"], font_name=FONT_BODY,
                 align=PP_ALIGN.LEFT)
    # Header underline
    add_rect(slide, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.02), C["neon_blue"])


def add_footer(slide, label="Mighty Skill-Bridge · 2026-06-02 CEO Brief"):
    """Bottom footer with brand label."""
    add_text(slide, label, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.3),
             font_size=9, color=C["text_muted"], font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)


def add_cta_box(slide, question, x=Inches(0.5), y=Inches(6.1), w=Inches(12.3), h=Inches(0.75)):
    """Highlighted question box for 社長への質問."""
    add_rect(slide, x, y, w, h, C["panel"], C["neon_blue"], Pt(1.5))
    # Q label
    add_text(slide, "▶ 社長への質問", Inches(0.75), y + Emu(50000), Inches(2.5), Inches(0.3),
             font_size=10, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Question text
    add_text(slide, question, Inches(0.75), y + Inches(0.32), w - Inches(0.5), Inches(0.42),
             font_size=13, color=C["text_primary"], bold=True, font_name=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_bullets(slide, points, x, y, w, h, *,
                font_size=14, color=None, bullet_color=None):
    color = color or C["text_secondary"]
    bullet_color = bullet_color or C["neon_green"]
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for idx, point in enumerate(points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if idx > 0:
            p.space_before = Pt(8)
        # Bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "■ "
        run_bullet.font.name = FONT_MONO
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        # Body
        run_body = p.add_run()
        run_body.text = point
        run_body.font.name = FONT_BODY
        run_body.font.size = Pt(font_size)
        run_body.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, *, label=None, accent=None):
    """Decorated panel (dark navy bg with left neon stripe)."""
    accent = accent or C["neon_blue"]
    add_rect(slide, x, y, w, h, C["panel"])
    add_rect(slide, x, y, Inches(0.08), h, accent)
    if label:
        add_text(slide, label, x + Inches(0.25), y + Inches(0.15),
                 w - Inches(0.4), Inches(0.3),
                 font_size=10, color=accent, bold=True, font_name=FONT_MONO,
                 align=PP_ALIGN.LEFT)


SLIDES = [
    {
        "num": 1,
        "title": "本日決めたいこと",
        "subtitle": "方針・優先順位・次アクションの決定 — 90 分想定",
        "points": [
            "実際の企画・サービス内容の最終決定ではなく、今後の方向性と優先順位を決める場",
            "現在までのプロトタイプと開発基盤の「実際にやった状態」を共有",
            "決定事項を打ち合わせ直後に WBS / Sheets / Calendar へ即時反映",
        ],
        "evidence": "本日のアジェンダ短文 (PRESHARE_MEMO 当日アジェンダ)",
        "question": "本日のゴールとして「方針決定と次回アクションの明確化」を設定してよろしいでしょうか?",
        "accent": "neon_blue",
    },
    {
        "num": 2,
        "title": "現在の到達点と公開デモ",
        "subtitle": "Seedance シネマティック UI 完成",
        "points": [
            "縦統合型シネマティックダッシュボードへ刷新、Seedance API 生成ブランドループ動画 + 非同期下部動画",
            "動画生成の非同期 polling + 外部 API 課金ガード (/admin ダッシュボード)",
            "AI 制限時 deterministic fallback + Public Demo Guard + Favicon / Chrome DevTools 対応完了",
        ],
        "evidence": "https://kanta13jp1.github.io/mighty-link-ai-connect/",
        "image_path": str(PROJECT_ROOT / "exports" / "knowledge_flow" / "screenshots" / "public_demo_hero_1920x1080.png"),
        "question": "最初の対象業務 (採用 / 営業 / SES 案件配属) のどれを最優先にしますか?",
        "accent": "neon_green",
    },
    {
        "num": 3,
        "title": "Google Workspace で進捗が回る基盤",
        "subtitle": "5 タブ Sheets + 完了済自動削除 Calendar",
        "points": [
            "WBS / Summary / Gantt 風 Timeline / 課題管理表 / QA 表 の計 5 タブを一括自動同期",
            "Google Calendar は 完了済イベントを自動削除 し未完了タスクだけが視覚的に残る",
            "OAuth アカウント検証 (k-umezawa@ml-mightylink.com) でセキュアな API 連携を固定",
        ],
        "evidence": "Sheets ID 1L99HCBHr4IsVUWqnUuG6OgoUmxEQUdfaYQim1n6etB8 · Calendar Mighty Skill-Bridge 開発計画",
        "question": "Google Workspace 連携は社内運用のみか、顧客提供価値の一部とするか?",
        "accent": "neon_blue",
    },
    {
        "num": 4,
        "title": "開発ナレッジ連携の実績とデモ",
        "subtitle": "NotebookLM / Notion / GitHub / Slack / Obsidian",
        "points": [
            "NotebookLM で資料要約 + 本プレゼンの構成案・PPTX 自動生成 (Drive 共有済)",
            "Notion 証跡ページ + GitHub Issues #1-#18 でタスク管理",
            "Slack 投稿案 + Obsidian vault 雛形の連携成果物生成",
        ],
        "evidence": "PPTX (Drive): https://docs.google.com/presentation/d/1XGHnQHBpJyyhh_Y3I2lq2UThPRC-2dcL    Notion: 3671d736b9db818aaa33da0a5f1a3951",
        "question": "NotebookLM / Slack / Notion / Obsidian のうち、6/2 以降の正式採用優先順位は?",
        "accent": "neon_green",
    },
    {
        "num": 5,
        "title": "サービス方向性の選択肢",
        "subtitle": "A / B / C のどれを育てるか (即決困難なら D: 保留可)",
        "points": [
            "方向性 A: AI フィット診断支援 — 営業 / 人材担当 / エンジニア向け",
            "方向性 B: Workspace 連携型 PM 支援 — 経営 / PM / 現場責任者向け",
            "方向性 C: AI PoC 高速構築支援 — 新規事業 / 営業企画 / 開発責任者向け",
        ],
        "evidence": "判断マトリクス: docs/CEO_PRESENTATION_DECISION_PACK_2026-06-02.md",
        "question": "このプロトタイプを何のサービスとして育て、最初に見せる相手は社内 / 既存顧客 / 見込み顧客のどれにしますか?",
        "accent": "neon_blue",
    },
    {
        "num": 6,
        "title": "運用・リスクの論点と公開範囲",
        "subtitle": "セキュリティ / 個人情報 / 共有範囲",
        "points": [
            "公開 URL の許容範囲 (社長共有 / 社内 / 既存顧客 / 外部) と認証層追加の要否",
            "個人情報 (経歴書等) の取り扱いと法務確認の時期 (方向性 A 選択時)",
            "Slack / Notion へ流す情報の共有範囲と権限マトリクス",
        ],
        "evidence": "リスク登録: data/issues_tracker.tsv (R9 個人情報, R10 公開 URL, R11 月額コスト)",
        "question": "公開 URL は社長共有 / 社内 / 外部公開のどこまで許容しますか? 法務確認時期は?",
        "accent": "neon_red",
    },
    {
        "num": 7,
        "title": "6/2 以降の優先開発機能と体制",
        "subtitle": "最初の 2 週間で何を作るか + 3-tool 体制コスト",
        "points": [
            "最優先機能 (スコア根拠強化 / 案件ストック管理 / WBS 内製化 等) を 1 つ決定",
            "3-tool 並走 (Antigravity / Codex / Claude Code) の API 月額コスト上限",
            "寛太の関与度 + 顧客 3 社同時パイロット時の追加採用シナリオ",
        ],
        "evidence": "Phase 7 ロードマップ枠: docs/CEO_PRESENTATION_POST_DECISION_ROADMAP_2026-06-02.md",
        "question": "最優先機能はどれにしますか? 開発用 AI の月額コスト上限は ¥10,000 / ¥30,000 / ¥50,000 のどれ?",
        "accent": "neon_green",
    },
    {
        "num": 8,
        "title": "次アクションと WBS への即時反映",
        "subtitle": "議事録 → Phase 7 WBS / Calendar / 課題管理表",
        "points": [
            "本日の決定事項 + 保留事項の整理 (議事録テンプレに即記入)",
            "Phase 7 WBS / Calendar / 課題管理表 / NotebookLM への即時反映を Codex レーンが実行",
            "次回レビュー日の設定 (推奨: 6/16 隔週 30 分定例化)",
        ],
        "evidence": "議事録テンプレ: CEO_PRESENTATION_DECISION_PACK · Notion CSV: notion_decision_log.csv",
        "question": "未決定で保留にしてよい事項は? 次回 (6/16 想定) で定例レビューにしてよいですか?",
        "accent": "neon_blue",
    },
]


def add_image_panel(slide, image_path, x, y, w, h, *, label="LIVE DEMO", caption=None):
    """Image panel: dark navy bg + neon left stripe + label header + screenshot + optional caption."""
    add_panel(slide, x, y, w, h, label=label, accent=C["neon_green"])
    img_x = x + Inches(0.2)
    img_y = y + Inches(0.5)
    img_w = w - Inches(0.4)
    img_h = h - (Inches(0.95) if caption else Inches(0.7))
    slide.shapes.add_picture(str(image_path), img_x, img_y, width=img_w, height=img_h)
    if caption:
        add_text(slide, caption,
                 x + Inches(0.2), y + h - Inches(0.4),
                 w - Inches(0.4), Inches(0.3),
                 font_size=9, color=C["text_secondary"], font_name=FONT_MONO,
                 align=PP_ALIGN.LEFT)


def render_one_screen_slide(prs, spec, gallery_idx, gallery_total):
    """One full-bleed screenshot per slide, with operation note + accent panel."""
    slide = base_slide(prs)
    # Composite header: "02.5-N / Demo Tour" + title
    num = f"02.{gallery_idx}"
    add_header(slide, num, spec["title"], spec.get("subtitle"))

    accent = C.get(spec.get("accent", "neon_green"), C["neon_green"])

    # Big screenshot centered horizontally; 16:9 ratio = 11.5" x 6.47" too tall
    # Use 10.5" x 5.91" — fits within 1.6 (header) + 5.91 + 0.95 (caption+CTA+footer) = 8.46 (need 7.5)
    # → smaller: 9.5" x 5.34", centered at x=1.92, y=1.6
    img_w = Inches(9.5)
    img_h = Inches(5.34)
    img_x = (SLIDE_W - img_w) / 2
    img_y = Inches(1.6)

    # Panel behind image
    add_rect(slide, img_x - Inches(0.1), img_y - Inches(0.1),
             img_w + Inches(0.2), img_h + Inches(0.2),
             C["panel"])
    # Top neon stripe
    add_rect(slide, img_x - Inches(0.1), img_y - Inches(0.1),
             img_w + Inches(0.2), Inches(0.05), accent)

    # Image
    img_path = Path(spec["path"])
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), img_x, img_y, width=img_w, height=img_h)

    # Right side label badge (operation context)
    if spec.get("op"):
        badge_x = img_x + img_w + Inches(0.1)
        badge_y = img_y
        # Right padding could overflow; clamp to slide
        badge_w = Inches(13.333) - badge_x - Inches(0.3)
        if badge_w > Inches(0.3):
            add_text(slide, spec["op"], badge_x, badge_y, badge_w, Inches(1.0),
                     font_size=10, bold=True, color=accent, font_name=FONT_MONO,
                     align=PP_ALIGN.LEFT)

    # Caption below image
    caption = spec.get("caption", "")
    cap_y = img_y + img_h + Inches(0.15)
    add_text(slide, caption, Inches(0.5), cap_y, Inches(12.3), Inches(0.45),
             font_size=14, color=C["text_secondary"], font_name=FONT_BODY,
             align=PP_ALIGN.CENTER)

    # Progress indicator dots (1 / 6 etc.)
    dots_y = Inches(6.75)
    dot_size = Inches(0.14)
    dot_gap = Inches(0.12)
    total_w = gallery_total * dot_size + (gallery_total - 1) * dot_gap
    dots_x = (SLIDE_W - total_w) / 2
    for i in range(gallery_total):
        dx = dots_x + i * (dot_size + dot_gap)
        color = accent if i == gallery_idx - 1 else C["text_muted"]
        add_rect(slide, dx, dots_y, dot_size, dot_size, color)

    add_footer(slide, f"Mighty Skill-Bridge · Demo Tour {gallery_idx}/{gallery_total} · 6/2 CEO Brief")


def render_slide(prs, spec):
    slide = base_slide(prs)
    add_header(slide, spec["num"], spec["title"], spec.get("subtitle"))

    accent = C.get(spec.get("accent", "neon_blue"), C["neon_blue"])

    # Left panel: KEY POINTS (60% width)
    add_panel(slide, Inches(0.5), Inches(1.6), Inches(8.0), Inches(4.0),
              label="KEY POINTS", accent=accent)
    add_bullets(slide, spec["points"],
                Inches(0.75), Inches(2.0), Inches(7.5), Inches(3.5),
                font_size=14, color=C["text_primary"], bullet_color=accent)

    # Right panel: IMAGE (if image_path) or EVIDENCE (text)
    img_path = spec.get("image_path")
    if img_path and Path(img_path).exists():
        add_image_panel(slide, img_path,
                        Inches(8.7), Inches(1.6), Inches(4.1), Inches(4.0),
                        label="LIVE DEMO",
                        caption=spec.get("evidence"))
    else:
        add_panel(slide, Inches(8.7), Inches(1.6), Inches(4.1), Inches(4.0),
                  label="EVIDENCE", accent=C["neon_green"])
        box = slide.shapes.add_textbox(Inches(8.9), Inches(2.0), Inches(3.7), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = spec["evidence"]
        run.font.name = FONT_MONO
        run.font.size = Pt(10)
        run.font.color.rgb = C["text_secondary"]

    # CTA: 社長への質問 (full width bottom band)
    add_cta_box(slide, spec["question"])
    add_footer(slide)


def render_title_slide(prs):
    """Title slide for the deck."""
    slide = base_slide(prs)
    # Background gradient effect via large panel
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5), C["panel"])
    # Cyan vertical line on left
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(0.12), Inches(3.5), C["neon_blue"])

    # Brand label
    add_text(slide, "MIGHTY SKILL-BRIDGE", Inches(0.8), Inches(2.2),
             Inches(11.5), Inches(0.4),
             font_size=14, color=C["neon_blue"], bold=True, font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)
    # Title
    add_text(slide, "CEO Brief 2026-06-02", Inches(0.8), Inches(2.7),
             Inches(11.5), Inches(0.9),
             font_size=44, color=C["text_primary"], bold=True, font_name=FONT_HEADER,
             align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(slide,
             "方針決定 + 優先順位 + 次アクション — 90 分で握る",
             Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
             font_size=18, color=C["neon_green"], font_name=FONT_BODY,
             align=PP_ALIGN.LEFT)
    # Tagline
    add_text(slide,
             "人とビジネスとテクノロジーを力強く繋ぐ — Mighty-Link Vision",
             Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
             font_size=12, color=C["text_secondary"], font_name=FONT_BODY,
             align=PP_ALIGN.LEFT)
    # Meta line
    add_text(slide,
             "Generated · 2026-05-23 · python-pptx · 8 slides · Workspace owned",
             Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.3),
             font_size=10, color=C["text_muted"], font_name=FONT_MONO,
             align=PP_ALIGN.LEFT)

    add_footer(slide, "Mighty Skill-Bridge · CEO Brief · 2026-06-02")


def build_summary(generated_at_jst):
    return f"""# Mighty Skill-Bridge CEO Presentation Deck (Branded)

Generated: {generated_at_jst.isoformat()}

## Output

- PPTX: `{PPTX_FILE.relative_to(PROJECT_ROOT)}`
- Generator: `scripts/generate_branded_ceo_deck.py`
- Source outline: `exports/knowledge_flow/notebooklm_ceo_slide_outline.md` (NotebookLM CLI 22 sources)
- Brand palette: Mighty Skill-Bridge (Seedance cinematic) — cyber black + neon blue + neon green

## Slides

1. Title slide (Brand cover)
2. 本日決めたいこと
3. 現在の到達点と公開デモ
4. Google Workspace で進捗が回る基盤
5. 開発ナレッジ連携の実績とデモ
6. サービス方向性の選択肢
7. 運用・リスクの論点と公開範囲
8. 6/2 以降の優先開発機能と体制
9. 次アクションと WBS への即時反映

## Design notes

- Cyber black background (#0d0e15) with neon blue/green/red accents
- Slide number badge + accent underline header
- 60/40 split: KEY POINTS panel (left) + EVIDENCE panel (right)
- Highlighted CTA box at bottom for 社長への質問
- Yu Gothic UI for Japanese body + Consolas for accents/IDs

## Companion

`mighty_skill_bridge_ceo_presentation_2026-06-02.pptx` (NotebookLM-CLI default style)
is the source-of-truth content. This `_branded.pptx` is the visual upgrade for
社長プレゼン 当日。Both decks have the same content; only styling differs.

## Next

- Run `python scripts/upload_notebooklm_docs_to_drive.py` to push to Drive (manual add to upload list if needed)
- Or use Canva MCP (HANDOFF-14b) for further refinement once OAuth is done
"""


def build_manifest(generated_at_jst):
    return {
        "generated_at_jst": generated_at_jst.isoformat(),
        "account": "k-umezawa@ml-mightylink.com",
        "generator": "scripts/generate_branded_ceo_deck.py",
        "source_outline": "exports/knowledge_flow/notebooklm_ceo_slide_outline.md",
        "source_count": 22,
        "brand_palette": {
            "bg": "#0D0E15",
            "panel": "#161824",
            "neon_blue": "#00F0FF",
            "neon_green": "#39FF14",
            "neon_red": "#FF3366",
            "text_primary": "#F1F5FF",
            "text_secondary": "#C5CAE0",
        },
        "outputs": {
            "pptx": str(PPTX_FILE.relative_to(PROJECT_ROOT)).replace("\\", "/"),
            "summary": str(SUMMARY_FILE.relative_to(PROJECT_ROOT)).replace("\\", "/"),
        },
        "slide_count": 9,
        "companion_deck": "exports/knowledge_flow/mighty_skill_bridge_ceo_presentation_2026-06-02.pptx",
        "purpose": "Canva/Figma template-style branded variant; visual upgrade only, same content as companion deck",
    }


SCREENSHOT_DIR = EXPORT_DIR / "screenshots"

GALLERY_CAPTURES = [
    {
        "title": "Step 1: 入力フォーム (Clean / 空欄)",
        "subtitle": "サイト初訪時 — 経歴書 + 案件票がプレースホルダー表示",
        "path": SCREENSHOT_DIR / "01_public_hero.png",
        "accent": "neon_green",
        "op": "操作: サイトにアクセスした直後 (auto-load なし)",
        "caption": "Mighty Skill-Bridge のエントリーポイント。サンプル自動投入を廃止し、本物の Step 1 として「Load Sample」または手入力を体験させる UX へ刷新。",
    },
    {
        "title": "Step 1: サンプルロード後の状態",
        "subtitle": "Load Sample で経歴書 + 案件票がプリフィル",
        "path": SCREENSHOT_DIR / "02_inputs_loaded.png",
        "accent": "neon_blue",
        "op": "操作: 「Load Sample」x2 押下",
        "caption": "実データに近いサンプルが投入された状態。社長デモではこの状態から Analyze を実行する。",
    },
    {
        "title": "Step 2: 4 軸フィット分析レポート",
        "subtitle": "Analyze 実行後 — deterministic fallback で構造化結果",
        "path": SCREENSHOT_DIR / "03_report_results.png",
        "accent": "neon_green",
        "op": "操作: 「Analyze Fit & Generate Story」押下",
        "caption": "Skill / Culture / Growth / Performing の 4 軸スコア + 面接質問生成。Gemini quota 中でも結果が出ることを実証。",
    },
    {
        "title": "Step 3: Seedance AI 動画デモ",
        "subtitle": "ブランドループ動画 + 非同期生成 + ダウンロード導線",
        "path": SCREENSHOT_DIR / "04_video_demo.png",
        "accent": "neon_blue",
        "op": "操作: 動画セクションへスクロール",
        "caption": "Seedance API で生成した Mighty Skill-Bridge ブランドビデオ。動画生成は非同期 polling + 課金ガード付き。",
    },
    {
        "title": "ナレッジ連携アーティファクト",
        "subtitle": "NotebookLM / Slack / Notion / Obsidian 成果物",
        "path": SCREENSHOT_DIR / "05_knowledge_flow.png",
        "accent": "neon_green",
        "op": "操作: Knowledge Flow セクション + Refresh",
        "caption": "Notebook Flow / Skill Bridge / Seedance / API Guard カードと、生成済 NotebookLM/Slack/Notion/Obsidian 成果物の状態確認。",
    },
    {
        "title": "管理者ダッシュボード /admin",
        "subtitle": "外部 API 課金ガード + Circuit Breaker + Recent Events",
        "path": SCREENSHOT_DIR / "06_admin_dashboard.png",
        "accent": "neon_red",
        "op": "操作: http://127.0.0.1:8000/admin",
        "caption": "Seedance + Gemini の日次呼び出し上限、サーキットブレーカー状態、ログイベントを 1 画面で監視。社内限定運用。",
    },
]


def main():
    generated_at = jst_now()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    render_title_slide(prs)

    # Content slides: insert per-screen gallery after slide 2 (現在の到達点と公開デモ)
    gallery_total = len(GALLERY_CAPTURES)
    for spec in SLIDES:
        render_slide(prs, spec)
        if spec.get("num") == 2:
            for idx, cap in enumerate(GALLERY_CAPTURES, start=1):
                render_one_screen_slide(prs, cap, idx, gallery_total)

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_FILE))
    SUMMARY_FILE.write_text(build_summary(generated_at), encoding="utf-8")
    MANIFEST_FILE.write_text(
        json.dumps(build_manifest(generated_at), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"[+] Branded CEO deck generated.")
    print(f"[*] PPTX:     {PPTX_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary:  {SUMMARY_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Manifest: {MANIFEST_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Slides:   {1 + len(SLIDES)} (1 title + {len(SLIDES)} content)")


if __name__ == "__main__":
    main()
