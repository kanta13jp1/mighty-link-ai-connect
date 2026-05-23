#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate a 10-pattern wireframe PPTX deck for the 6/2 CEO presentation.

Greyscale wireframe aesthetic (no brand colors) so the conversation is about
"layout pattern" rather than visual design. Each slide presents one wireframe
variant of the Mighty Skill-Bridge AI Fit Diagnosis service with:
- Header: pattern number, title, subtitle
- Body: layout illustration (rectangles, dashed boxes, pills, placeholder lines)
- Footer: design rationale (1-line, why this pattern matters)

Companion to the Figma file aiQt3c1Cenru4x6GMcLuL5 (same 10 patterns, but
Figma file is empty due to Starter plan MCP rate limit at creation time).

Output:
- exports/knowledge_flow/mighty_skill_bridge_ui_wireframes_2026-06-02.pptx
- exports/knowledge_flow/mighty_skill_bridge_ui_wireframes_2026-06-02.md
- exports/knowledge_flow/mighty_skill_bridge_ui_wireframes_2026-06-02.json
"""

from __future__ import annotations

import datetime as dt
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = PROJECT_ROOT / "exports" / "knowledge_flow"
PPTX_FILE = EXPORT_DIR / "mighty_skill_bridge_ui_wireframes_2026-06-02.pptx"
SUMMARY_FILE = EXPORT_DIR / "mighty_skill_bridge_ui_wireframes_2026-06-02.md"
MANIFEST_FILE = EXPORT_DIR / "mighty_skill_bridge_ui_wireframes_2026-06-02.json"

# 16:9 1920x1080 equivalent
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Greyscale wireframe palette
C = {
    "bg":       RGBColor(0xFA, 0xFB, 0xFC),
    "panel":    RGBColor(0xF0, 0xF2, 0xF5),
    "panel2":   RGBColor(0xE6, 0xE9, 0xEE),
    "stroke":   RGBColor(0xD0, 0xD5, 0xDB),
    "text_p":   RGBColor(0x1F, 0x29, 0x37),
    "text_s":   RGBColor(0x6B, 0x72, 0x80),
    "text_m":   RGBColor(0x9C, 0xA3, 0xAF),
    "cta":      RGBColor(0x37, 0x41, 0x51),
    "cta_text": RGBColor(0xFA, 0xFB, 0xFC),
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_BODY = "Yu Gothic UI"
FONT_MONO = "Consolas"


def jst_now() -> dt.datetime:
    return dt.datetime.now(dt.timezone(dt.timedelta(hours=9)))


def add_rect(slide, x, y, w, h, fill_color=None, *, stroke=None, stroke_w=None, dashed=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        if stroke_w is not None:
            shape.line.width = stroke_w
        if dashed:
            try:
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            except Exception:
                pass
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, *, stroke=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h, *, size=11, bold=False,
             color=None, font_name=None, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    color = color or C["text_p"]
    font_name = font_name or FONT_BODY
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def base_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C["bg"])
    return slide


def header(slide, num, title, subtitle):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.7), C["panel"])
    add_text(slide, f"WF-{num:02d}", Inches(0.25), Inches(0.18), Inches(0.8), Inches(0.3),
             size=11, bold=True, color=C["text_m"], font_name=FONT_MONO)
    add_text(slide, title, Inches(1.1), Inches(0.13), Inches(11.5), Inches(0.32),
             size=18, bold=True, color=C["text_p"])
    add_text(slide, subtitle, Inches(1.1), Inches(0.42), Inches(11.5), Inches(0.24),
             size=11, color=C["text_s"])
    add_rect(slide, 0, Inches(0.71), SLIDE_W, Inches(0.012), C["stroke"])


def footer(slide, rationale):
    add_rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.012), C["stroke"])
    add_text(slide, "💡 " + rationale, Inches(0.25), Inches(7.12), Inches(13.0), Inches(0.3),
             size=11, bold=True, color=C["text_s"])


def dashed_box(slide, x, y, w, h, label=None):
    add_rect(slide, x, y, w, h, C["white"], stroke=C["text_m"], stroke_w=Pt(0.75), dashed=True)
    if label:
        add_text(slide, label, x, y, w, h,
                 size=10, color=C["text_m"], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def pill(slide, x, y, w, h, label, filled=False):
    bg = C["cta"] if filled else C["white"]
    fg = C["cta_text"] if filled else C["text_p"]
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if not filled:
        shape.line.color.rgb = C["stroke"]
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.5
    shape.shadow.inherit = False
    add_text(slide, label, x, y, w, h,
             size=10, bold=True, color=fg, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def placeholder_lines(slide, x, y, w, count, gap=Inches(0.12)):
    line_h = Inches(0.07)
    for i in range(count):
        line_w = w * 0.6 if i == count - 1 else w
        add_rect(slide, x, y + i * (line_h + gap), line_w, line_h, C["panel2"])


# === 10 wireframe builders ===

def wf01(s):
    """Vertical Hero Stack — Mobile-first"""
    dashed_box(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(1.3),
               "▶ Hero animation / brand statement")
    add_text(s, "1. 経歴書をアップロード", Inches(0.5), Inches(2.4), Inches(6.0), Inches(0.3),
             size=14, bold=True)
    dashed_box(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.7),
               "📄 Drag & drop here / Click to browse")
    add_text(s, "2. 案件票を貼り付け", Inches(0.5), Inches(3.6), Inches(6.0), Inches(0.3),
             size=14, bold=True)
    dashed_box(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.7),
               "📋 Paste JD text or attach file")
    pill(s, Inches(5.0), Inches(4.85), Inches(3.3), Inches(0.55),
         "Analyze & Generate Story", filled=True)
    add_text(s, "3. 結果", Inches(0.5), Inches(5.55), Inches(4.0), Inches(0.25),
             size=12, bold=True, color=C["text_s"])
    placeholder_lines(s, Inches(0.5), Inches(5.85), Inches(12.3), 4)


def wf02(s):
    """Split Form — Profile vs Job (Current)"""
    col_w = Inches(5.85)
    # left col
    add_rounded_rect(s, Inches(0.5), Inches(0.95), col_w, Inches(5.0), C["panel"])
    add_text(s, "Engineer · Profile", Inches(0.7), Inches(1.1), col_w - Inches(0.4),
             Inches(0.25), size=14, bold=True)
    placeholder_lines(s, Inches(0.7), Inches(1.5), col_w - Inches(0.4), 3)
    dashed_box(s, Inches(0.7), Inches(2.1), col_w - Inches(0.4), Inches(3.0),
               "textarea: 経歴・スキル・希望条件")
    pill(s, Inches(0.7), Inches(5.25), Inches(1.4), Inches(0.32), "Load Sample")
    # right col
    add_rounded_rect(s, Inches(6.95), Inches(0.95), col_w, Inches(5.0), C["panel"])
    add_text(s, "Job · Requirements", Inches(7.15), Inches(1.1), col_w - Inches(0.4),
             Inches(0.25), size=14, bold=True)
    placeholder_lines(s, Inches(7.15), Inches(1.5), col_w - Inches(0.4), 3)
    dashed_box(s, Inches(7.15), Inches(2.1), col_w - Inches(0.4), Inches(3.0),
               "textarea: 必須・歓迎要件・カルチャー")
    pill(s, Inches(7.15), Inches(5.25), Inches(1.4), Inches(0.32), "Load Sample")
    # CTA
    pill(s, Inches(5.0), Inches(6.15), Inches(3.3), Inches(0.55),
         "Analyze Fit & Generate Story", filled=True)


def wf03(s):
    """Step Wizard — 4-step"""
    steps = ["1. Profile", "2. Job", "3. Review", "4. Results"]
    step_w = (Inches(13.333) - Inches(0.8)) / 4
    for i in range(4):
        sx = Inches(0.4) + i * step_w
        active = i == 1
        # progress bar segment
        bar_color = C["cta"] if i < 2 else C["panel2"]
        add_rect(s, sx, Inches(1.05), step_w - Inches(0.15), Inches(0.04), bar_color)
        # node circle
        node_color = C["cta"] if active else C["white"]
        node_x = sx + (step_w - Inches(0.15)) / 2 - Inches(0.18)
        add_rounded_rect(s, node_x, Inches(0.95), Inches(0.36), Inches(0.36),
                         node_color, stroke=C["stroke"])
        add_text(s, str(i + 1), node_x, Inches(0.95), Inches(0.36), Inches(0.36),
                 size=11, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 color=C["white"] if active else C["text_s"])
        add_text(s, steps[i], sx, Inches(1.4), step_w - Inches(0.15), Inches(0.25),
                 size=11, bold=True, align=PP_ALIGN.CENTER,
                 color=C["text_p"] if active else C["text_s"])
    # active panel
    add_rounded_rect(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.5),
                     C["panel"])
    add_text(s, "Step 2 · 案件票を入力", Inches(0.75), Inches(2.05), Inches(11.0),
             Inches(0.25), size=14, bold=True)
    dashed_box(s, Inches(0.75), Inches(2.45), Inches(11.8), Inches(3.8),
               "textarea: 案件詳細...")
    pill(s, Inches(0.75), Inches(6.55), Inches(1.4), Inches(0.36), "← Back")
    pill(s, Inches(11.4), Inches(6.55), Inches(1.4), Inches(0.36), "Next →", filled=True)


def wf04(s):
    """Conversational Chat"""
    # bot bubble 1
    add_rounded_rect(s, Inches(0.5), Inches(0.95), Inches(6.0), Inches(0.55), C["panel"])
    add_text(s, "🤖 AI: 経歴を教えてください。職種、経験年数、得意領域。",
             Inches(0.65), Inches(0.95), Inches(5.7), Inches(0.55),
             size=11, anchor=MSO_ANCHOR.MIDDLE)
    # user bubble 1
    add_rounded_rect(s, Inches(6.83), Inches(1.6), Inches(6.0), Inches(0.55), C["cta"])
    add_text(s, "Python 8 年 / FastAPI / 採用系 SaaS の MVP 設計",
             Inches(6.98), Inches(1.6), Inches(5.7), Inches(0.55),
             size=11, color=C["cta_text"], anchor=MSO_ANCHOR.MIDDLE)
    # bot bubble 2
    add_rounded_rect(s, Inches(0.5), Inches(2.25), Inches(6.0), Inches(0.55), C["panel"])
    add_text(s, "🤖 AI: 次に、希望案件の規模感は?",
             Inches(0.65), Inches(2.25), Inches(5.7), Inches(0.55),
             size=11, anchor=MSO_ANCHOR.MIDDLE)
    # user typing
    add_rounded_rect(s, Inches(6.83), Inches(2.9), Inches(6.0), Inches(0.55),
                     C["white"], stroke=C["stroke"])
    add_text(s, "_ Type your reply...",
             Inches(6.98), Inches(2.9), Inches(5.7), Inches(0.55),
             size=11, color=C["text_m"], anchor=MSO_ANCHOR.MIDDLE)
    # result preview
    add_rounded_rect(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.5), C["panel2"])
    add_text(s, "💎 暫定 Fit Score", Inches(0.7), Inches(3.85), Inches(6.0), Inches(0.25),
             size=12, bold=True)
    add_text(s, "Skill 72 · Culture — · Growth — · Performing —",
             Inches(0.7), Inches(4.15), Inches(12.0), Inches(0.25),
             size=11, color=C["text_s"])
    placeholder_lines(s, Inches(0.7), Inches(4.5), Inches(11.8), 3, gap=Inches(0.07))
    # input area
    add_rounded_rect(s, Inches(0.5), Inches(5.5), Inches(11.0), Inches(0.45),
                     C["white"], stroke=C["stroke"])
    add_text(s, "Type message...",
             Inches(0.7), Inches(5.5), Inches(10.8), Inches(0.45),
             size=11, color=C["text_m"], anchor=MSO_ANCHOR.MIDDLE)
    pill(s, Inches(11.7), Inches(5.5), Inches(1.1), Inches(0.45), "Send", filled=True)


def wf05(s):
    """Drag & Drop Cards — 4 Quadrant"""
    labels = ["Engineers", "Jobs", "Matches", "Reports"]
    half_w = Inches(6.0)
    row_h = Inches(2.8)
    for i in range(4):
        qx = Inches(0.5) + (i % 2) * (half_w + Inches(0.3))
        qy = Inches(0.95) + (i // 2) * (row_h + Inches(0.2))
        add_rounded_rect(s, qx, qy, half_w, row_h, C["panel"])
        add_text(s, labels[i], qx + Inches(0.2), qy + Inches(0.15),
                 half_w - Inches(0.4), Inches(0.25), size=13, bold=True)
        for k in range(3):
            cy = qy + Inches(0.5) + k * Inches(0.7)
            add_rounded_rect(s, qx + Inches(0.2), cy,
                             half_w - Inches(0.4), Inches(0.6),
                             C["white"], stroke=C["stroke"])
            add_text(s, f"{labels[i][:-1]} #{k + 1}",
                     qx + Inches(0.35), cy + Inches(0.1),
                     half_w - Inches(0.8), Inches(0.2),
                     size=11, bold=True)
            add_text(s, f"Fit: 8{3 - k}",
                     qx + half_w - Inches(1.1), cy + Inches(0.1),
                     Inches(0.9), Inches(0.2),
                     size=11, color=C["text_s"], align=PP_ALIGN.RIGHT)


def wf06(s):
    """Pipeline Board — Kanban"""
    cols = ["候補", "評価中", "推薦", "配属済"]
    counts = [12, 5, 3, 8]
    col_w = (Inches(13.333) - Inches(1.0) - Inches(0.45)) / 4
    col_h = Inches(5.4)
    for i in range(4):
        cx = Inches(0.5) + i * (col_w + Inches(0.15))
        add_rounded_rect(s, cx, Inches(0.95), col_w, col_h, C["panel"])
        add_text(s, cols[i], cx + Inches(0.15), Inches(1.1),
                 col_w - Inches(0.7), Inches(0.25), size=13, bold=True)
        add_text(s, f"({counts[i]})", cx + col_w - Inches(0.6), Inches(1.1),
                 Inches(0.5), Inches(0.25), size=11, color=C["text_s"],
                 align=PP_ALIGN.RIGHT)
        for k in range(4):
            cy = Inches(1.5) + k * Inches(0.92)
            add_rounded_rect(s, cx + Inches(0.15), cy,
                             col_w - Inches(0.3), Inches(0.78),
                             C["white"], stroke=C["stroke"])
            add_text(s, f"候補 {(i * 4) + k + 1}",
                     cx + Inches(0.28), cy + Inches(0.08),
                     col_w - Inches(0.9), Inches(0.2),
                     size=11, bold=True)
            add_text(s, f"案件 × Skill 8{(9 - k - i) % 10}",
                     cx + Inches(0.28), cy + Inches(0.3),
                     col_w - Inches(0.9), Inches(0.2),
                     size=10, color=C["text_s"])
            # score badge
            badge_x = cx + col_w - Inches(0.55)
            add_rounded_rect(s, badge_x, cy + Inches(0.1),
                             Inches(0.36), Inches(0.22), C["cta"])
            add_text(s, f"A{(9 - i) % 10}", badge_x, cy + Inches(0.1),
                     Inches(0.36), Inches(0.22),
                     size=9, bold=True, color=C["cta_text"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def wf07(s):
    """Comparison Table — Matrix"""
    cols = ["Name", "Skill", "Culture", "Growth", "Performing", "Total"]
    col_w = (Inches(13.333) - Inches(1.0)) / len(cols)
    # header
    add_rounded_rect(s, Inches(0.5), Inches(0.95), Inches(13.333) - Inches(1.0),
                     Inches(0.42), C["panel"])
    for i, lbl in enumerate(cols):
        add_text(s, lbl, Inches(0.5) + i * col_w + Inches(0.15),
                 Inches(0.95), col_w - Inches(0.3), Inches(0.42),
                 size=12, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # rows
    rows = [
        ("田中 A", 92, 85, 78, 88, "85.8"),
        ("佐藤 B", 88, 90, 82, 86, "86.5"),
        ("鈴木 C", 76, 82, 95, 80, "83.3"),
        ("高橋 D", 95, 70, 72, 90, "81.8"),
        ("伊藤 E", 80, 88, 84, 85, "84.3"),
        ("渡辺 F", 70, 80, 90, 75, "78.8"),
        ("山本 G", 85, 92, 88, 82, "86.8"),
    ]
    for ri, row in enumerate(rows):
        ry = Inches(1.45) + ri * Inches(0.6)
        add_rect(s, Inches(0.5), ry, Inches(13.333) - Inches(1.0), Inches(0.55),
                 C["white"] if ri % 2 == 0 else C["bg"],
                 stroke=C["stroke"], stroke_w=Pt(0.5))
        for ci, val in enumerate(row):
            tx = Inches(0.5) + ci * col_w + Inches(0.15)
            add_text(s, str(val), tx, ry + Inches(0.05),
                     col_w - Inches(0.3), Inches(0.2),
                     size=11, bold=(ci == 0 or ci == len(row) - 1))
            if isinstance(val, int):
                ratio = val / 100
                bar_w = (col_w - Inches(0.3)) * ratio
                bar_color = C["cta"] if val >= 85 else C["text_m"]
                add_rect(s, tx, ry + Inches(0.32), bar_w, Inches(0.08), bar_color)


def wf08(s):
    """Dashboard Tiles"""
    titles = [
        ("本日のマッチ", "+12"),
        ("Score 分布", "📊"),
        ("Top Skill Gap", "Python"),
        ("Recent Activity", "8 events"),
        ("Interview Queue", "3 pending"),
        ("Monthly Hires", "+5"),
    ]
    tile_w = (Inches(13.333) - Inches(1.0) - Inches(0.3)) / 3
    tile_h = Inches(2.7)
    for i, info in enumerate(titles):
        tx = Inches(0.5) + (i % 3) * (tile_w + Inches(0.15))
        ty = Inches(0.95) + (i // 3) * (tile_h + Inches(0.2))
        add_rounded_rect(s, tx, ty, tile_w, tile_h, C["panel"])
        add_text(s, info[0], tx + Inches(0.2), ty + Inches(0.15),
                 tile_w - Inches(0.4), Inches(0.22),
                 size=11, color=C["text_s"])
        add_text(s, info[1], tx + Inches(0.2), ty + Inches(0.4),
                 tile_w - Inches(0.4), Inches(0.7),
                 size=32, bold=True)
        dashed_box(s, tx + Inches(0.2), ty + Inches(1.3),
                   tile_w - Inches(0.4), tile_h - Inches(1.5),
                   "📈 sparkline / chart placeholder")


def wf09(s):
    """Inline Preview — Live"""
    left_w = Inches(7.5)
    right_w = Inches(13.333) - left_w - Inches(0.9)
    # left
    add_rounded_rect(s, Inches(0.5), Inches(0.95), left_w - Inches(0.2),
                     Inches(5.7), C["panel"])
    add_text(s, "経歴書 + 案件票 (入力中)", Inches(0.7), Inches(1.1),
             left_w - Inches(0.6), Inches(0.25), size=14, bold=True)
    dashed_box(s, Inches(0.7), Inches(1.5), left_w - Inches(0.6),
               Inches(4.9), "live typing area — 入力するたびに右が更新")
    # right
    add_rounded_rect(s, Inches(0.5) + left_w, Inches(0.95), right_w,
                     Inches(5.7), C["white"], stroke=C["stroke"])
    add_text(s, "🔴 LIVE Fit Preview",
             Inches(0.7) + left_w, Inches(1.1),
             right_w - Inches(0.4), Inches(0.25),
             size=13, bold=True)
    axes = ["Skill", "Culture", "Growth", "Performing"]
    scores = [78, 65, 82, 71]
    for i, (a, sc) in enumerate(zip(axes, scores)):
        ay = Inches(1.6) + i * Inches(0.55)
        add_text(s, a, Inches(0.7) + left_w, ay, Inches(0.9), Inches(0.22),
                 size=11, bold=True)
        add_rect(s, Inches(1.6) + left_w, ay + Inches(0.05),
                 right_w - Inches(2.1), Inches(0.1), C["panel2"])
        add_rect(s, Inches(1.6) + left_w, ay + Inches(0.05),
                 (right_w - Inches(2.1)) * (sc / 100),
                 Inches(0.1), C["cta"])
        add_text(s, str(sc),
                 Inches(0.5) + left_w + right_w - Inches(0.6), ay,
                 Inches(0.4), Inches(0.22),
                 size=11, bold=True, align=PP_ALIGN.RIGHT)
    # candidates
    add_text(s, "近い候補 (top 3)",
             Inches(0.7) + left_w, Inches(4.1),
             right_w - Inches(0.4), Inches(0.22),
             size=11, bold=True, color=C["text_s"])
    for k in range(3):
        cy = Inches(4.4) + k * Inches(0.55)
        add_rounded_rect(s, Inches(0.7) + left_w, cy,
                         right_w - Inches(0.4), Inches(0.45),
                         C["bg"], stroke=C["stroke"])
        add_text(s, f"候補 {k + 1}",
                 Inches(0.85) + left_w, cy + Inches(0.1),
                 Inches(1.5), Inches(0.25),
                 size=11, bold=True)
        add_text(s, f"Skill 8{4 - k}",
                 Inches(0.5) + left_w + right_w - Inches(1.0),
                 cy + Inches(0.1), Inches(0.85), Inches(0.25),
                 size=11, color=C["text_s"], align=PP_ALIGN.RIGHT)


def wf10(s):
    """Search-First Catalog"""
    # search
    add_rounded_rect(s, Inches(0.5), Inches(0.95), Inches(10.5), Inches(0.5),
                     C["white"], stroke=C["stroke"])
    add_text(s, "🔍 検索: スキル / 業界 / 経験年数...",
             Inches(0.7), Inches(0.95), Inches(10.0), Inches(0.5),
             size=12, color=C["text_m"], anchor=MSO_ANCHOR.MIDDLE)
    pill(s, Inches(11.2), Inches(0.95), Inches(1.6), Inches(0.5),
         "+ New Match", filled=True)
    # filter chips
    chips = [("Python", True), ("Senior", False), ("Remote", False),
             ("JP", False), ("Hourly", False), ("+ Add", False)]
    cx = Inches(0.5)
    for label, filled in chips:
        w_pill = Inches(0.55) if label == "+ Add" else Inches(0.85)
        pill(s, cx, Inches(1.65), w_pill, Inches(0.32), label, filled=filled)
        cx += w_pill + Inches(0.1)
    # grid
    card_w = (Inches(13.333) - Inches(1.0) - Inches(0.3)) / 3
    card_h = Inches(2.1)
    for i in range(6):
        xx = Inches(0.5) + (i % 3) * (card_w + Inches(0.15))
        yy = Inches(2.15) + (i // 3) * (card_h + Inches(0.15))
        add_rounded_rect(s, xx, yy, card_w, card_h, C["panel"])
        # avatar
        add_rounded_rect(s, xx + Inches(0.15), yy + Inches(0.15),
                         Inches(0.48), Inches(0.48),
                         C["white"], stroke=C["stroke"])
        add_text(s, f"候補 #{i + 1}",
                 xx + Inches(0.75), yy + Inches(0.2),
                 card_w - Inches(0.9), Inches(0.22),
                 size=13, bold=True)
        add_text(s, "Python · FastAPI · 8y",
                 xx + Inches(0.75), yy + Inches(0.42),
                 card_w - Inches(0.9), Inches(0.22),
                 size=11, color=C["text_s"])
        add_rect(s, xx + Inches(0.15), yy + Inches(0.95),
                 card_w - Inches(0.3), Inches(0.1), C["panel2"])
        add_rect(s, xx + Inches(0.15), yy + Inches(0.95),
                 (card_w - Inches(0.3)) * 0.85, Inches(0.1), C["cta"])
        add_text(s, "Fit 85",
                 xx + Inches(0.15), yy + Inches(1.1),
                 card_w - Inches(0.3), Inches(0.22),
                 size=11, color=C["text_s"])
        pill(s, xx + Inches(0.15), yy + Inches(1.45),
             Inches(0.7), Inches(0.28), "Remote")
        pill(s, xx + Inches(0.9), yy + Inches(1.45),
             Inches(0.7), Inches(0.28), "Senior")


VARIANTS = [
    {"num": 1, "title": "Vertical Hero Stack — Mobile-first",
     "subtitle": "1 column / top-to-bottom flow. 縦スクロールで完結。スマホ親和。",
     "rationale": "最小 UI で迷わず Step 1→2→3 を順に体験できる。SES 営業が客先で iPad プレゼン時に強い。",
     "build": wf01},
    {"num": 2, "title": "Split Form — Profile vs Job (Current)",
     "subtitle": "2 列並列入力 / 中央 Analyze CTA。現行 UI のベース。",
     "rationale": "比較対象を視覚的に対置。デスクトップでの広い画面前提。",
     "build": wf02},
    {"num": 3, "title": "Step Wizard — 4-step Progress Flow",
     "subtitle": "経歴 → 案件 → 確認 → 結果。1 ステップずつ集中。",
     "rationale": "初めて使う社内人事担当が迷わない。各ステップで Help が出せる。",
     "build": wf03},
    {"num": 4, "title": "Conversational Chat — AI Interview",
     "subtitle": "AI が質問 / ユーザーが回答 / 最終 fit を結果バブルで提示。",
     "rationale": "応募者本人が直接使う場合に親しみやすい。チャット履歴で説明責任。",
     "build": wf04},
    {"num": 5, "title": "Drag & Drop Cards — 4 Quadrant Workspace",
     "subtitle": "Engineers / Jobs / Matches / Reports の 4 象限。カード DnD。",
     "rationale": "中規模の人事チームが複数案件 × 複数候補を同時管理。",
     "build": wf05},
    {"num": 6, "title": "Pipeline Board — Kanban for Recruiters",
     "subtitle": "候補 → 評価中 → 推薦 → 配属済。Score badge 付きカード。",
     "rationale": "リクルーター業務の標準ビュー。複数案件並列追跡しやすい。",
     "build": wf06},
    {"num": 7, "title": "Comparison Table — Multi-candidate Matrix",
     "subtitle": "候補 (rows) × 評価軸 (cols)。color-coded score。",
     "rationale": "経営報告/最終決裁で最強。1 画面で 5-10 候補を比較できる。",
     "build": wf07},
    {"num": 8, "title": "Dashboard Tiles — Exec Overview",
     "subtitle": "6 タイル: 本日マッチ / Score 分布 / Skill Gap / Recent / Queue / Metrics。",
     "rationale": "社長 / 経営層が一目で全体把握。週次レビューで使う。",
     "build": wf08},
    {"num": 9, "title": "Inline Preview — Live Suggestions",
     "subtitle": "入力中に右パネルへリアルタイム fit / 候補をストリーミング表示。",
     "rationale": "AI ライブ感を最大化。社長デモで「打ちながら結果が変わる」を見せる。",
     "build": wf09},
    {"num": 10, "title": "Search-First Catalog — Engineer/Job Library",
     "subtitle": "検索バー + フィルタ → グリッドカード → 詳細モーダル。",
     "rationale": "蓄積データが増えた後の \"リクルーターの日常入口\"。検索駆動。",
     "build": wf10},
]


def render_title_slide(prs):
    slide = base_slide(prs)
    add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3),
                     Inches(3.5), C["panel"])
    add_text(slide, "WIREFRAME COLLECTION", Inches(0.8), Inches(2.2),
             Inches(11.5), Inches(0.4),
             size=12, bold=True, color=C["text_s"], font_name=FONT_MONO)
    add_text(slide, "Mighty Skill-Bridge UI Patterns", Inches(0.8),
             Inches(2.65), Inches(11.5), Inches(0.9),
             size=44, bold=True, color=C["text_p"])
    add_text(slide, "10 design pattern variants — for CEO direction selection",
             Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
             size=18, color=C["text_s"])
    add_text(slide, "Greyscale wireframes · 16:9 · Generated 2026-05-24",
             Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
             size=12, color=C["text_m"], font_name=FONT_MONO)
    add_text(slide, "Pair file (Figma): https://www.figma.com/design/aiQt3c1Cenru4x6GMcLuL5",
             Inches(0.8), Inches(4.85), Inches(11.5), Inches(0.4),
             size=10, color=C["text_m"], font_name=FONT_MONO)


def main():
    generated_at = jst_now()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    render_title_slide(prs)

    for v in VARIANTS:
        slide = base_slide(prs)
        header(slide, v["num"], v["title"], v["subtitle"])
        v["build"](slide)
        footer(slide, v["rationale"])

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_FILE))

    summary = f"""# Mighty Skill-Bridge UI Wireframes — 10 Patterns

Generated: {generated_at.isoformat()}

## Output

- PPTX: `{PPTX_FILE.relative_to(PROJECT_ROOT).as_posix()}`
- Slides: 1 title + 10 wireframes = 11 slides
- Style: greyscale wireframe (no brand colors — focus on layout pattern)

## Companion

- Figma file (empty container, awaiting MCP rate limit reset):
  https://www.figma.com/design/aiQt3c1Cenru4x6GMcLuL5

## 10 Patterns

""" + "\n".join(
        f"{v['num']:02d}. **{v['title']}** — {v['subtitle']}\n    _Rationale_: {v['rationale']}"
        for v in VARIANTS
    ) + "\n"
    SUMMARY_FILE.write_text(summary, encoding="utf-8")

    manifest = {
        "generated_at_jst": generated_at.isoformat(),
        "account": "k-umezawa@ml-mightylink.com",
        "generator": "scripts/generate_wireframes_deck.py",
        "outputs": {
            "pptx": PPTX_FILE.relative_to(PROJECT_ROOT).as_posix(),
            "summary": SUMMARY_FILE.relative_to(PROJECT_ROOT).as_posix(),
        },
        "slide_count": 1 + len(VARIANTS),
        "style": "greyscale wireframe",
        "figma_companion_file_key": "aiQt3c1Cenru4x6GMcLuL5",
        "figma_companion_url": "https://www.figma.com/design/aiQt3c1Cenru4x6GMcLuL5",
        "patterns": [
            {"num": v["num"], "title": v["title"], "subtitle": v["subtitle"],
             "rationale": v["rationale"]}
            for v in VARIANTS
        ],
    }
    MANIFEST_FILE.write_text(json.dumps(manifest, ensure_ascii=False, indent=2),
                             encoding="utf-8")

    print(f"[+] Wireframes deck generated.")
    print(f"[*] PPTX:     {PPTX_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Summary:  {SUMMARY_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Manifest: {MANIFEST_FILE.relative_to(PROJECT_ROOT)}")
    print(f"[*] Slides:   {1 + len(VARIANTS)} (1 title + {len(VARIANTS)} wireframes)")


if __name__ == "__main__":
    main()
